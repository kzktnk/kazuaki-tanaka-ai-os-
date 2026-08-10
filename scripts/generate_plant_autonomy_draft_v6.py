#!/usr/bin/env python3
"""Generate plant autonomy draft v6 — definition-first, consolidated (9 slides)."""

import sys
from pathlib import Path

_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from _config import output_paths

OUT_PATHS = output_paths("plant-autonomy_operating-model_draft_v6.pptx")

IBM_BLUE = RGBColor(0x05, 0x3F, 0x87)
IBM_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
ORANGE = RGBColor(0xC4, 0x50, 0x00)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_LINE = RGBColor(0x2E, 0x7D, 0x32)
TEXT = RGBColor(0x33, 0x33, 0x33)
FOOTER = "クライアント様 発電所自立経営 定義・論点整理（たたき台 v6）｜2026年8月7日"


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_section_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.20), Inches(12.70), Inches(0.61))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_key_message(slide, text):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.85), Inches(12.70), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_footer(slide, section, num):
    f = slide.shapes.add_textbox(Inches(0.50), Inches(7.18), Inches(10.50), Inches(0.25))
    f.text_frame.text = f"{FOOTER}　｜　{section}"
    f.text_frame.paragraphs[0].font.size = Pt(8)
    f.text_frame.paragraphs[0].font.color.rgb = GRAY
    n = slide.shapes.add_textbox(Inches(12.60), Inches(7.18), Inches(0.40), Inches(0.25))
    n.text_frame.text = str(num)
    n.text_frame.paragraphs[0].font.size = Pt(8)
    n.text_frame.paragraphs[0].font.color.rgb = GRAY
    n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, font_size=11):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, bold = item if isinstance(item, tuple) else (item, False)
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.font.bold = bold
        p.space_after = Pt(4)


def add_round_box(slide, left, top, width, height, fill, line, items, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    add_bullets(slide, left + Inches(0.20), top + Inches(0.10), width - Inches(0.40), height - Inches(0.20), items, font_size)
    return shape


def add_table(slide, data, left, top, width, height, font_size=10):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = IBM_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True


def build_title_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    t1 = slide.shapes.add_textbox(Inches(0.32), Inches(0.94), Inches(12.70), Inches(2.81))
    t1.text_frame.text = "クライアント様\n\n発電所自立経営に向けた\n「定義」と「論点整理」（たたき台）"
    for p in t1.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE
    t2 = slide.shapes.add_textbox(Inches(0.44), Inches(4.24), Inches(6.04), Inches(0.63))
    t2.text_frame.text = "2026年8月7日"
    t2.text_frame.paragraphs[0].font.size = Pt(14)
    t3 = slide.shapes.add_textbox(Inches(9.27), Inches(5.68), Inches(2.71), Inches(0.63))
    t3.text_frame.text = "日本アイ・ビー・エム株式会社"
    t3.text_frame.paragraphs[0].font.size = Pt(12)
    add_footer(slide, "表紙", slide_no)


def build_purpose_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "本資料の目的と構成")
    add_key_message(
        slide,
        "体制案の提示ではなく、まず「自立運営モデル発電所とは何か」を言語化し、"
        "それに付随する論点を整理することが目的",
    )
    add_round_box(
        slide,
        Inches(0.50),
        Inches(1.75),
        Inches(5.90),
        Inches(2.20),
        IBM_LIGHT,
        IBM_BLUE,
        [
            ("本資料の位置づけ", True),
            "・IBMの仮説（たたき台）。クライアント様とすり合わせながら議論する前提",
            "・エグゼクティブサマリーは本編確定後に別途作成（本資料には含まない）",
            "・体制図・権限レベル詳細は、定義合意後に具体化",
        ],
    )
    add_round_box(
        slide,
        Inches(6.70),
        Inches(1.75),
        Inches(6.13),
        Inches(2.20),
        GREEN_BG,
        GREEN_LINE,
        [
            ("構成（キーメッセージでストーリーを追える）", True),
            "1. なぜ今、言語化が先か",
            "2. 自立運営モデル発電所とは何か",
            "3. 3つの取組の位置付け",
            "4. 何を実証するか → 5. How論点 → 6. 次のステップ",
        ],
    )
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(4.20), Inches(12.33), Inches(1.05))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    note.line.color.rgb = ORANGE
    add_bullets(
        slide,
        Inches(0.70),
        Inches(4.30),
        Inches(11.90),
        Inches(0.85),
        [
            ("経営層ディスカッションのねらい（案）", True),
            "・「自立運営モデル発電所」の定義に合意する　／　How論点（権限・KPI・関与・期間）の優先順位を決める",
        ],
        11,
    )
    add_footer(slide, "目的と構成", slide_no)


def build_why_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "1. なぜ今、言語化が先か")
    add_key_message(
        slide,
        "4月の経営責任移管だけでは「自立経営」は始まらない。"
        "責任・権限・実行力をセットで決める必要がある",
    )
    add_round_box(
        slide,
        Inches(0.50),
        Inches(1.70),
        Inches(6.00),
        Inches(2.45),
        IBM_LIGHT,
        IBM_BLUE,
        [
            ("いま起きていること", True),
            "・4月より発電所長が経営責任者に。DX戦略・AMプログラムでも発電所主体を目指している",
            "・一方で「何をどこまで任せるか」がまだ決まっていない",
            "・AI検証（トラックA）の議論が進んでも、自立運営の前提が共有されていない",
        ],
    )
    add_round_box(
        slide,
        Inches(6.70),
        Inches(1.70),
        Inches(6.13),
        Inches(2.45),
        GREEN_BG,
        GREEN_LINE,
        [
            ("IBMの考え（仮説）", True),
            "・足りないのは体制図ではなく「自立運営モデル発電所」の定義",
            "・責任だけ移しても、判断・投資・改善の権限がなければ経営はできない",
            "・定義が固まってから、権限・KPI・体制・期間を具体化する",
        ],
    )
    add_round_box(
        slide,
        Inches(0.50),
        Inches(4.35),
        Inches(12.33),
        Inches(1.20),
        RGBColor(0xFF, 0xF3, 0xE0),
        ORANGE,
        [
            ("So What", True),
            "・本資料は「How（体制案）」の前に「What（定義）」を議論するためのたたき台",
            "・ギャップ整理・体制案・権限レベル詳細は、定義合意後の具体化フェーズで扱う",
        ],
        11,
    )
    add_footer(slide, "Why", slide_no)


def build_what_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "2. 自立運営モデル発電所とは何か（IBM仮説）")
    add_key_message(
        slide,
        "AIを試す場所ではなく、発電所が経営の主体になる仕組みを1拠点で試す場",
    )
    add_round_box(
        slide,
        Inches(0.50),
        Inches(1.65),
        Inches(6.00),
        Inches(2.05),
        IBM_LIGHT,
        IBM_BLUE,
        [
            ("定義", True),
            "・発電所長がKPI達成に向け、任された範囲で自ら判断・実行するモデル",
            "・KPIの管理は本社、現場の実行権限は発電所——この二層で回す",
            "・AMプログラム：収支計画PDCA　／　DX：デジタル施策の実行手段",
        ],
        10,
    )
    add_round_box(
        slide,
        Inches(6.70),
        Inches(1.65),
        Inches(6.13),
        Inches(2.05),
        GREEN_BG,
        GREEN_LINE,
        [
            ("IBMの考え", True),
            "・AIは目的ではなく、経営力を高める道具",
            "・モデル発電所＝AI導入実験ではなく、運営の仕組みの実証",
            "・同時に試すもの：①経営 ②組織 ③PDCA ④AI活用（AIは4番目）",
        ],
        10,
    )
    data = [
        ["", "トラックA（AI検証）", "トラックB 自立運営モデル実証"],
        ["目的", "技術・導入の可否確認", "権限×KPI×PDCAの運用モデル実証"],
        ["位置づけ", "トラックBの入力になりうる", "本資料の主題"],
        ["関係", "別トラックで並行", "定義合意後に具体設計"],
    ]
    add_table(slide, data, Inches(0.50), Inches(3.95), Inches(12.33), Inches(2.05), 10)
    add_bullets(
        slide,
        Inches(0.50),
        Inches(6.15),
        Inches(12.33),
        Inches(0.45),
        ["※ トラックA（AI検証）を進める前提として、まずトラックBの定義を固める必要がある"],
        9,
    )
    add_footer(slide, "What", slide_no)


def build_positioning_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "3. DX戦略 × アセット管理構築プログラム × モデルプラント（位置付け）")
    add_key_message(
        slide,
        "3つの取組は別プロジェクトではなく、運営の仕組みを変える1つの変革の異なるレイヤー",
    )

    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(1.55), Inches(12.33), Inches(0.35))
    banner.fill.solid()
    banner.fill.fore_color.rgb = IBM_LIGHT
    banner.line.color.rgb = IBM_BLUE
    add_bullets(slide, Inches(0.65), Inches(1.58), Inches(12.00), Inches(0.30), ["【上段】DX（AI活用）戦略"], 10)

    steps = [
        ("① 戦略策定", "重点施策抽出／推進体制整理"),
        ("② ロードマップ", "業務モデル設計／KPI設定／検証"),
        ("③ 展開", "全所展開／結果を①②へ還元"),
    ]
    x = 0.55
    for i, (title, body) in enumerate(steps):
        left = Inches(x)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.95), Inches(3.55), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = IBM_BLUE
        add_bullets(slide, left + Inches(0.12), Inches(2.00), Inches(3.30), Inches(0.85), [(title, True), body], 9)
        x += 3.95
        if i < 2:
            arr = slide.shapes.add_textbox(Inches(x - 0.35), Inches(2.20), Inches(0.30), Inches(0.40))
            arr.text_frame.text = "→"
            arr.text_frame.paragraphs[0].font.size = Pt(14)
            arr.text_frame.paragraphs[0].font.bold = True

    amp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(3.05), Inches(12.33), Inches(0.75))
    amp.fill.solid()
    amp.fill.fore_color.rgb = GREEN_BG
    amp.line.color.rgb = GREEN_LINE
    add_bullets(
        slide,
        Inches(0.70),
        Inches(3.12),
        Inches(11.90),
        Inches(0.60),
        [
            ("【横串】アセット管理構築プログラム", True),
            "発電所主体の収支計画PDCA・モニタリング｜2027年度事業計画プロセスへの実装着手が今年度ゴール",
        ],
        9,
    )

    add_bullets(slide, Inches(0.50), Inches(3.95), Inches(12.00), Inches(0.25), ["【下段】モデルプラント（2トラック）"], 10)

    add_round_box(
        slide,
        Inches(0.50),
        Inches(4.25),
        Inches(6.00),
        Inches(1.45),
        WHITE,
        ACCENT,
        [
            ("トラックA（AI検証）", True),
            "目的：技術・導入の可否確認（PoC）",
            "時期：2026/9〜2027/3（案）",
        ],
        9,
    )
    add_round_box(
        slide,
        Inches(6.80),
        Inches(4.25),
        Inches(6.03),
        Inches(1.45),
        WHITE,
        ORANGE,
        [
            ("トラックB 自立運営モデル実証", True),
            "目的：権限×KPI×PDCAの運用モデル実証",
            "時期：2026/10〜2027/3（案）｜本資料の主題",
        ],
        9,
    )

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(5.85), Inches(12.33), Inches(0.55))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    note.line.color.rgb = ORANGE
    add_bullets(
        slide,
        Inches(0.70),
        Inches(5.92),
        Inches(11.90),
        Inches(0.40),
        [
            "So What：モデル発電所は「AI検証」ではなく、新しい運営の仕組みを試す場。"
            "トラックAとトラックBは目的が異なるため分離して設計する",
        ],
        9,
    )
    add_footer(slide, "位置付け", slide_no)


def build_prove_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "4. モデル発電所で何を実証するか")
    add_key_message(slide, "主役はAIではなく「経営の仕組み」。AIはその一部（手段）")
    data = [
        ["優先順", "実証する内容", "補足"],
        ["1", "意思決定（誰が何を決めるか）", "権限の範囲と例外ルール"],
        ["2", "投資・改善の判断", "予算・施策選定の裁量"],
        ["3", "データ活用とPDCA", "AMプログラムの収支計画サイクルと接続"],
        ["4", "KPI管理と現場実行の両立", "本社の目標設定と現場の達成責任"],
        ["5", "AIの活用", "技術は変わる。運営モデルは長く使う"],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.70), Inches(12.33), Inches(3.35), 10)
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(5.25), Inches(12.33), Inches(1.05))
    note.fill.solid()
    note.fill.fore_color.rgb = IBM_LIGHT
    note.line.color.rgb = ACCENT
    add_bullets(
        slide,
        Inches(0.70),
        Inches(5.35),
        Inches(11.90),
        Inches(0.85),
        [
            ("PoCとの違い", True),
            "・技術検証（PoC）ではなく、1拠点・一定期間の「運営モデル実証」",
            "・成功基準は技術評価だけでなく、KPI達成・定着・権限設計の妥当性を含む",
        ],
        10,
    )
    add_footer(slide, "実証内容", slide_no)


def build_how_topics_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "5. 定義のあとに議論すべき論点（How）")
    add_key_message(
        slide,
        "体制・権限・KPI・期間は、定義合意後に順番に具体化する。"
        "先に体制案を出すと本末転倒",
    )
    boxes = [
        (
            Inches(0.50),
            Inches(1.65),
            "① 権限",
            "・ヒト・モノ・カネ・情報をどこまで任せるか\n"
            "・社内ルール・マニュアルの例外（治外法権）はどのレベルまで許容するか",
        ),
        (
            Inches(6.70),
            Inches(1.65),
            "② 本社・他発電所の関与",
            "・モデル発電所の実現に、本社は何を支援・統制するか\n"
            "・他発電所はいつ、どう横展開に関わるか",
        ),
        (
            Inches(0.50),
            Inches(3.55),
            "③ KPI",
            "・ストレッチ目標を誰が設定するか\n"
            "・どの会議体で約束（コミット）するか／未達時の介入ルール",
        ),
        (
            Inches(6.70),
            Inches(3.55),
            "④ 期間",
            "・技術検証は短期（数ヶ月）でも可\n"
            "・経営成果の確認には最低1経営サイクル（例：6〜12ヶ月）が必要",
        ),
    ]
    for left, top, title, body in boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.00), Inches(1.70)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = IBM_LIGHT
        shape.line.color.rgb = IBM_BLUE
        add_bullets(
            slide,
            left + Inches(0.20),
            top + Inches(0.12),
            Inches(5.60),
            Inches(1.45),
            [(title, True), body],
            10,
        )
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(5.45), Inches(12.33), Inches(0.95))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    note.line.color.rgb = ORANGE
    add_bullets(
        slide,
        Inches(0.70),
        Inches(5.55),
        Inches(11.90),
        Inches(0.75),
        [
            ("体制について", True),
            "・PMO・二層ライン・権限レベル（Level 0–3）等は、上記4論点の合意後に設計する",
            "・詳細な体制案・ギャップ表・Optionsは次フェーズの具体化資料で扱う",
        ],
        10,
    )
    add_footer(slide, "How論点", slide_no)


def build_next_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "6. IBMの考え方と次のステップ")
    add_key_message(
        slide,
        "発電所自立経営＝権限委譲でもAI導入でもなく、"
        "発電所が価値を創り続ける運営の仕組みへの転換",
    )
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(1.65), Inches(12.33), Inches(1.05))
    quote.fill.solid()
    quote.fill.fore_color.rgb = IBM_BLUE
    quote.line.color.rgb = IBM_BLUE
    box = slide.shapes.add_textbox(Inches(0.80), Inches(1.85), Inches(11.70), Inches(0.75))
    tf = box.text_frame
    tf.text = (
        "IBMの主張：発電所自立経営とは、権限委譲でもAI導入でもない。\n"
        "発電所が継続的に価値を創出できる「運営の仕組み（Operating Model）」への転換である。"
    )
    for p in tf.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
    data = [
        ["時期", "アクション", "成果物・ねらい"],
        ["すり合わせ期間", "PO・プロジェクトオーナーとのすり合わせ", "定義・論点のたたき台合意"],
        ["経営層ディスカッション", "経営層ディスカッション", "定義合意 ＋ How論点4つの優先順位"],
        ["9月以降", "トラックB具体設計", "拠点選定・権限/KPI設計・体制案・スケジュール"],
        ["来週", "エグゼクティブサマリー作成", "本編確定後に別途提示"],
    ]
    add_table(slide, data, Inches(0.50), Inches(2.95), Inches(12.33), Inches(2.55), 10)
    add_footer(slide, "次のステップ", slide_no)


def build_eof_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.20), Inches(12.70), Inches(0.61))
    box.text_frame.text = "End of Document"
    box.text_frame.paragraphs[0].font.size = Pt(24)
    box.text_frame.paragraphs[0].font.bold = True
    add_footer(slide, "End", slide_no)


def main():
    dst = Presentation()
    dst.slide_width = Inches(13.33)
    dst.slide_height = Inches(7.50)

    build_title_slide(dst, 1)
    build_purpose_slide(dst, 2)
    build_why_slide(dst, 3)
    build_what_slide(dst, 4)
    build_positioning_slide(dst, 5)
    build_prove_slide(dst, 6)
    build_how_topics_slide(dst, 7)
    build_next_slide(dst, 8)
    build_eof_slide(dst, 9)

    for out in OUT_PATHS:
        out.parent.mkdir(parents=True, exist_ok=True)
        dst.save(str(out))
        print(f"Saved: {out}")
    print(f"Total slides: {len(dst.slides)}")


if __name__ == "__main__":
    main()
