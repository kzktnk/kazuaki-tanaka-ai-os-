#!/usr/bin/env python3
"""Generate plant autonomy operating-model draft PPTX v3."""

import sys
from pathlib import Path

_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from _config import output_paths

OUT_PATHS = output_paths("plant-autonomy_operating-model_draft_v3.pptx")
OUT = str(OUT_PATHS[0])

IBM_BLUE = RGBColor(0x05, 0x3F, 0x87)
IBM_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
ORANGE = RGBColor(0xC4, 0x50, 0x00)


def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title, subtitle=None, slide_no=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = IBM_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.58), Inches(9.4), Inches(0.35))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(10)
        stf.paragraphs[0].font.color.rgb = GRAY
    if slide_no:
        num = slide.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
        ntf = num.text_frame
        ntf.text = str(slide_no)
        ntf.paragraphs[0].font.size = Pt(9)
        ntf.paragraphs[0].font.color.rgb = GRAY
        ntf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullet_box(slide, left, top, width, height, bullets, font_size=11):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level, bold = item[0], item[1] if len(item) > 1 else 0, item[2] if len(item) > 2 else False
        else:
            text, level, bold = item, 0, False
        p.text = text
        p.level = level
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.bold = bold
        p.space_after = Pt(4)
    return box


def add_table(slide, rows, cols, left, top, width, height, data, header=True):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = IBM_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.size = Pt(9)
    return table


def slide_title(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    t1.text_frame.text = "骨子（たたき台）Rev.3"
    t1.text_frame.paragraphs[0].font.size = Pt(14)
    t1.text_frame.paragraphs[0].font.color.rgb = GRAY

    t2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(1.2))
    tf = t2.text_frame
    tf.text = "発電所自立経営に向けた\n「体制」と「モデルプラント」の設計"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = IBM_BLUE

    t3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.5))
    tf3 = t3.text_frame
    tf3.text = (
        "① 推進体制への示唆と体制案（案）\n"
        "② デジタル施策実行に係る自立運営モデルプラントの設計（案）"
    )
    for p in tf3.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    meta = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
    mtf = meta.text_frame
    mtf.text = (
        "クライアント様　技術統括部 デジタルパワープラント推進部\n"
        "PO・ステークホルダー　ご報告用ドラフト\n"
        "2026年8月5日（第2回ヒアリング）時点の情報に基づく仮説ベースの一次たたき台　｜　IBM"
    )
    for p in mtf.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
    return slide


def slide_purpose(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "本資料の位置づけ", slide_no=n)

    add_bullet_box(slide, Inches(0.4), Inches(0.95), Inches(4.5), Inches(2.2), [
        ("本日（8/5）のご要望", 0, True),
        ("① 推進体制に対する示唆だし（改善提案）と、貴社側で実際の体制案を記述", 0),
        ("② デジタル施策実行に係る自立運営モデルプラントの設計（案）を記述", 0),
    ], 10)

    add_bullet_box(slide, Inches(0.4), Inches(3.2), Inches(4.5), Inches(2.0), [
        ("Rev.3での主な更新", 0, True),
        ("・3つの取組（DX戦略／AMP自立経営／モデルプラント）の統合プロセス図を追加", 0),
        ("・As-Is／Gap／To-Beによる現状ギャップ整理を追加", 0),
        ("・Rev.4形式の体制図・権限マトリクス（Decision Ownership 3層）を具体化", 0),
        ("・自立運営モデルプラントの検証設計・経営判断Optionsを追加", 0),
    ], 10)

    add_bullet_box(slide, Inches(5.0), Inches(0.95), Inches(4.5), Inches(2.5), [
        ("活用タイミング", 0, True),
        ("8/7（金）　簡易インプットとしてご提供", 0),
        ("8/12　定例会議でのPOへのご報告", 0),
        ("8/13以降　経営層ディスカッション（権限設定・自立経営モデル）", 0),
    ], 10)

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.5), Inches(4.5), Inches(1.5))
    note.fill.solid()
    note.fill.fore_color.rgb = IBM_LIGHT
    note.line.color.rgb = ACCENT
    ntf = note.text_frame
    ntf.text = (
        "ご留意：本資料は仮説ベースの骨子です。"
        "用語はクライアント資料に合わせ「発電所自立経営」に統一しています。"
        "8/12以降の議論・フィードバックを経て精緻化する前提でご覧ください。"
    )
    for p in ntf.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY


def slide_process(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "全体プロセス｜3つの取組の関係性", "DX（AI活用）推進 × AMP構築PJ × モデルプラント", n)

    # Upper band
    upper = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0), Inches(9.2), Inches(1.5))
    upper.fill.solid()
    upper.fill.fore_color.rgb = IBM_LIGHT
    upper.line.color.rgb = IBM_BLUE
    utf = upper.text_frame
    utf.text = "【上段】DX（AI活用）戦略の策定・実行プロセス"
    utf.paragraphs[0].font.bold = True
    utf.paragraphs[0].font.size = Pt(11)
    utf.paragraphs[0].font.color.rgb = IBM_BLUE

    boxes = [
        ("① 戦略策定\nDGD施策抽出\n重点施策選定\n推進体制整理", Inches(0.6), Inches(1.45)),
        ("→", Inches(2.55), Inches(1.65)),
        ("② ロードマップ具体化\n業務モデル設計\nKPI設定\n開発・検証", Inches(2.85), Inches(1.45)),
        ("→", Inches(4.8), Inches(1.65)),
        ("展開\nフィードバック\n全発電所へ", Inches(5.1), Inches(1.45)),
    ]
    for text, left, top in boxes:
        if text == "→":
            arr = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.3))
            arr.text_frame.text = "→"
            arr.text_frame.paragraphs[0].font.size = Pt(16)
            arr.text_frame.paragraphs[0].font.bold = True
        else:
            b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.9), Inches(0.95))
            b.fill.solid()
            b.fill.fore_color.rgb = WHITE
            b.line.color.rgb = IBM_BLUE
            btf = b.text_frame
            btf.text = text
            for p in btf.paragraphs:
                p.font.size = Pt(8)
                p.alignment = PP_ALIGN.CENTER

    # Middle AMP
    mid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(2.7), Inches(9.2), Inches(1.0))
    mid.fill.solid()
    mid.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    mid.line.color.rgb = ORANGE
    mtf = mid.text_frame
    mtf.text = (
        "【横串】AMP構築PJ｜発電所自立経営・アセマネプロセス構築\n"
        "発電所主体の収支計画PDCA・モニタリング｜2027年度事業計画プロセス（2027年）への実装着手が今年度ゴール"
    )
    for p in mtf.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = ORANGE
        if p == mtf.paragraphs[0]:
            p.font.bold = True

    # Lower model plants
    lower_label = slide.shapes.add_textbox(Inches(0.4), Inches(3.85), Inches(9), Inches(0.3))
    lower_label.text_frame.text = "【下段】モデルプラント（2トラック）— 目的が異なるため分離して設計"
    lower_label.text_frame.paragraphs[0].font.bold = True
    lower_label.text_frame.paragraphs[0].font.size = Pt(10)
    lower_label.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE

    mp1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(4.3), Inches(1.5))
    mp1.fill.solid()
    mp1.fill.fore_color.rgb = WHITE
    mp1.line.color.rgb = ACCENT
    mp1tf = mp1.text_frame
    mp1tf.text = (
        "③-a AI検証モデルプラント（既存・Rev.4）\n"
        "目的：クイック施策・PoCの技術・導入検証\n"
        "時期：2026/9〜2027/3\n"
        "主体：発電所長＋機能チーム（DPP/G-DAC）"
    )
    for p in mp1tf.paragraphs:
        p.font.size = Pt(9)

    mp2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(4.2), Inches(4.5), Inches(1.5))
    mp2.fill.solid()
    mp2.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    mp2.line.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
    mp2.line.width = Pt(2)
    mp2tf = mp2.text_frame
    mp2tf.text = (
        "③-b 自立運営モデルプラント（新設・本資料の提案）\n"
        "目的：権限委譲×KPI管理の運用モデル実証\n"
        "時期：2026/10〜2027/3（案）\n"
        "主体：発電所長＋PMO拡張機能＋AMP連携"
    )
    for p in mp2tf.paragraphs:
        p.font.size = Pt(9)
        if p == mp2tf.paragraphs[0]:
            p.font.bold = True

    add_bullet_box(slide, Inches(0.4), Inches(5.85), Inches(9.2), Inches(1.0), [
        ("So What：上段②と横串AMPは「合流していく」前提。③-aと③-bは別トラックで並行推進し、③-bで経営判断に必要な権限設計の実証データを得る", 0, True),
    ], 9)


def slide_as_is_gap(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "As-Is → Gap → To-Be", "4月移行後の実態ギャップ整理（Pattern 2）", n)

    data = [
        ["", "As-Is（現状）", "Gap（構造的課題）", "To-Be（あるべき姿）"],
        [
            "組織・\n権限",
            "4月より発電所長が経営責任者に。\nCJPOから数値目標を受領済み",
            "ツール選定・予算執行等の\n権限範囲が未明文化",
            "権限委譲マトリクスに基づく\n現場裁量＋本社ガードレール",
        ],
        [
            "KPI\n管理",
            "Rev.4で機能⑥に発電所が主担。\nCJPOが目標設定",
            "KPI管理（本社集中）と\n実行権限（現場委譲）の非対称",
            "KPI管理ラインと実行権限ラインを\n分離し接続点を明確化",
        ],
        [
            "プロセス",
            "本社主導のシステム導入が\n従来の標準パターン",
            "現場主導導入時の\n失敗・定着リスクへの備え不足",
            "Plan→Build→Run→Improve\n＋エスカレーション支援",
        ],
        [
            "モデル\nプラント",
            "AI検証用モデルプラント\n（③-a）のみ計画",
            "自立経営モデル（③-b）の\n検証場が体制上未整備",
            "2トラック並行＋AMP PDCA\nを題材とした実証",
        ],
        [
            "PJ\n連携",
            "DX推進体制とAMP構築PJが\n別体系で整理",
            "全体最適と現場自律の\n関係性が未整理",
            "PMO機能共通化・\n定期報告統合",
        ],
    ]
    add_table(slide, 6, 4, Inches(0.25), Inches(0.95), Inches(9.5), Inches(5.8), data)


def slide_why(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Why｜なぜ今、体制とモデルプラントの設計が必要か", slide_no=n)

    add_bullet_box(slide, Inches(0.4), Inches(0.95), Inches(4.5), Inches(2.8), [
        ("現状の到達点", 0, True),
        ("・DX（AI活用）推進体制（案）はRev.4で策定済み", 0),
        ("・AMP構築PJで「発電所自律運営によるユニット価値最大化」を目指す姿を設定", 0),
        ("・4月から発電所長が経営責任者に移行、CJPOから数値目標を受領", 0),
        ("・今年度ゴール：2027年度事業計画プロセスへのPDCA実装着手", 0),
    ], 10)

    add_bullet_box(slide, Inches(0.4), Inches(3.9), Inches(4.5), Inches(2.8), [
        ("足元の経営文脈", 0, True),
        ("・事業領域毎の戦略が不透明な中でDX戦略を策定する必要（PO指摘）", 0),
        ("・8/13以降、経営層と権限設定・自立経営モデルのあり方を議論", 0),
    ], 10)

    add_bullet_box(slide, Inches(5.0), Inches(0.95), Inches(4.5), Inches(5.5), [
        ("未解決の論点＝示唆の源泉", 0, True),
        ("1. 権限委譲（現場裁量）とKPI管理（本社集中）の非対称性を制度設計として両立", 0),
        ("2. モデルプラント＝クイックAI検証（③-a）であり、自立経営実証（③-b）の場が未整備", 0),
        ("3. DX推進体制とAMP構築PJの関係性が未整理（全体最適 vs 現場自律）", 0),
        ("4. 現場主導の施策導入における失敗・定着リスクへのガバナンス設計", 0),
        ("5. 発電所長が「困らないようサポート」する仕組みの体制への落とし込み", 0),
    ], 10)


def slide_rev4_comments(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "①体制への示唆｜Rev.4体制案に対するコメント", slide_no=n)

    data = [
        ["#", "観点", "現状（Rev.4）", "示唆", "So What"],
        ["1", "権限の明文化", "発電所長は意見提示・整合判断・先行検証にとどまる", "権限委譲マトリクスを体制表に追加", "現場の施策停滞を防ぐ"],
        ["2", "KPI/実行権限", "KPI管理が機能⑥・PMO・統括部長に分散", "KPI管理ラインと実行権限ラインを分離", "本社統制と現場自律の両立"],
        ["3", "エスカレーション", "KPI未達時フロー未記載", "発電所長→PMO→CJPOの支援ライン新設", "「困らないサポート」を制度化"],
        ["4", "PJ接続", "DX推進とAMP構築PJが別体系", "PMO共通化・定期報告統合", "合流前提の全体最適"],
        ["5", "モデルプラント", "AI検証用途に限定", "③-b自立運営モデル実証を別トラック新設", "経営判断の実証データ取得"],
        ["6", "会議体", "発電所長はモデルプラント時のみ必須", "全体進捗会議等への位置づけ明確化", "全所への展開意思醸成"],
        ["7", "機能⑥矛盾", "KPI主担=発電所 vs 本社集中方針", "KPI設定（本社）/達成責任（現場）/報告（PMO）を分離", "Rev.4内部矛盾の解消"],
    ]
    add_table(slide, 8, 5, Inches(0.15), Inches(0.9), Inches(9.7), Inches(5.9), data)


def slide_org_chart(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "①体制案（To-Be）｜全体像", "Rev.4形式｜KPI管理ライン（本社集中）と実行権限ライン（現場委譲）の二層構造", n)

    # CJPO
    cjpo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(0.95), Inches(3.0), Inches(0.55))
    cjpo.fill.solid()
    cjpo.fill.fore_color.rgb = IBM_BLUE
    cjpo.line.fill.background()
    cjpo.text_frame.text = "CJPO（CJPO担当）｜最終意思決定"
    cjpo.text_frame.paragraphs[0].font.color.rgb = WHITE
    cjpo.text_frame.paragraphs[0].font.size = Pt(10)
    cjpo.text_frame.paragraphs[0].font.bold = True
    cjpo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # PMO
    pmo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(1.7), Inches(5.0), Inches(0.7))
    pmo.fill.solid()
    pmo.fill.fore_color.rgb = ACCENT
    pmo.line.fill.background()
    pmo.text_frame.text = "PM/PMO（統括部長／PO・PMO）＋自律経営運営機能（PMO拡張）"
    pmo.text_frame.paragraphs[0].font.color.rgb = WHITE
    pmo.text_frame.paragraphs[0].font.size = Pt(9)
    pmo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    roles = [
        ("本社実行責任\n各統括部長", Inches(0.3), Inches(2.6), IBM_LIGHT, IBM_BLUE),
        ("発電所実行責任\n発電所長", Inches(3.5), Inches(2.6), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x2E, 0x7D, 0x32)),
        ("実行支援\nDPP推進部・G-DAC\nデジタル部門", Inches(6.8), Inches(2.6), IBM_LIGHT, IBM_BLUE),
    ]
    for text, left, top, fill, line in roles:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = line
        box.text_frame.text = text
        for p in box.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER

    # Two lines
    kpi_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(3.95), Inches(9.4), Inches(0.55))
    kpi_line.fill.solid()
    kpi_line.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    kpi_line.line.color.rgb = RGBColor(0xC6, 0x28, 0x28)
    kpi_line.text_frame.text = "KPI管理ライン（本社集中）：CJPO目標設定 → PMO横串モニタリング → 統括部長整合 → 発電所長達成責任"
    kpi_line.text_frame.paragraphs[0].font.size = Pt(9)
    kpi_line.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xB7, 0x1C, 0x1C)

    exec_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.6), Inches(9.4), Inches(0.55))
    exec_line.fill.solid()
    exec_line.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    exec_line.line.color.rgb = ACCENT
    exec_line.text_frame.text = "実行権限ライン（現場委譲）：発電所長が権限マトリクス範囲内で施策・ツール選定・予算執行｜機能チームがガードレール・技術助言"
    exec_line.text_frame.paragraphs[0].font.size = Pt(9)
    exec_line.text_frame.paragraphs[0].font.color.rgb = IBM_BLUE

    esc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(5.25), Inches(9.4), Inches(0.45))
    esc.fill.solid()
    esc.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    esc.line.color.rgb = ORANGE
    esc.text_frame.text = "エスカレーション（未達時）：発電所長 → PMO（支援・調整） → CJPO（重要判断・権限見直し）"
    esc.text_frame.paragraphs[0].font.size = Pt(9)

    add_bullet_box(slide, Inches(0.3), Inches(5.85), Inches(9.4), Inches(1.0), [
        ("注：「自律経営運営事務局」は新組織増設ではなく、PMO機能の拡張（AMP/DX接続窓口・権限運用ルール整備）として位置づけ", 0, True),
        ("Rev.4 機能①–⑥（業務要件・データ基盤・AI開発・遠隔監視・現場実装・KPI管理）との対応関係は別紙または次版で詳細化", 0),
    ], 9)


def slide_authority_matrix(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "①体制案｜権限委譲マトリクス（案）＋ Decision Ownership", slide_no=n)

    data = [
        ["領域", "委譲（発電所長）", "保持（本社/CJPO）", "Reviewer", "Approver", "Owner", "エスカレーション"],
        ["ツール選定", "カタログ内・予算内", "基盤・セキュリティ基準", "DPP/G-DAC", "発電所長", "発電所長", "カタログ外・大型投資"],
        ["予算執行", "KPI達成に資する範囲", "年度枠・重点配分", "PMO", "発電所長/CJPO", "発電所長", "予算超過見込"],
        ["KPI", "目標水準の提案", "設定・改定", "統括部長", "CJPO", "CJPO", "3ヶ月連続未達"],
        ["人員配置", "所内配置裁量", "横断配置方針", "統括部長", "CJPO", "発電所長", "大幅不足"],
        ["業務プロセス\n変更", "所内運用設計", "全社標準・BPR", "DPP", "統括部長", "発電所長", "全社標準逸脱"],
        ["データ整備", "所内データ入力・整備", "マスタ・ガバナンス", "DPP", "統括部長", "発電所長", "基盤不整合"],
    ]
    add_table(slide, 7, 7, Inches(0.1), Inches(0.9), Inches(9.8), Inches(3.8), data)

    add_bullet_box(slide, Inches(0.3), Inches(4.85), Inches(9.4), Inches(2.0), [
        ("Decision Ownership 設計原則（GitHub frameworks/decision-delegation 準拠）", 0, True),
        ("・実行権限は委譲可能。結果の説明責任（Accountability）はOwnerに残る", 0),
        ("・AI/デジタル施策：AI推薦 → 人間Review → 人間Approve → 実行 → Outcome Ownership", 0),
        ("・委譲ラインの定量基準（予算上限・期間等）は8/13経営層ディスカッションで確定（本表は論点整理用）", 0),
    ], 9)


def slide_amp_connection(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "AMP構築PJとの接続", "自立運営モデルプラント（③-b）をAMP PDCA実証の場として位置づけ", n)

    data = [
        ["AMP構築PJ要素", "現状/方向性", "③-bモデルプラントでの検証内容"],
        ["目指す姿", "発電所自律運営によるユニット価値最大化", "権限委譲下でのKPI達成→ユニット価値への寄与を実測"],
        ["今年度ゴール", "2027年度事業計画プロセスへのPDCA実装着手", "収支計画→実績→改善アクションのPDCAを1拠点で試行"],
        ["発電所長の役割", "社長型経営責任者（4月移行済）", "PDCAサイクルの自律運用＋デジタル施策選定"],
        ["PO所管", "AMシステム基盤構築部会リーダー", "KPI定義・データ要件の③-bへの接続"],
        ["D-61等施策", "ユニット別・発電所別収支管理", "ダッシュボード等を題材施策として活用"],
    ]
    add_table(slide, 6, 3, Inches(0.3), Inches(0.95), Inches(9.4), Inches(3.5), data)

    add_bullet_box(slide, Inches(0.3), Inches(4.6), Inches(9.4), Inches(2.2), [
        ("検証の題材施策（例示）", 0, True),
        ("A. 自立運営ダッシュボード（DPP AP：SAP業務基盤活用・5件構築予定）", 0),
        ("B. クイックWin施策（計画外停止削減系）— ③-a AI検証と連携可能", 0),
        ("C. 収支計画PDCAプロセス — AMP構築PJの2027事業計画実装の先行実証", 0),
    ], 10)


def slide_model_why_what(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "②自立運営モデルプラントの設計（案）｜Why / What", slide_no=n)

    add_bullet_box(slide, Inches(0.4), Inches(0.95), Inches(4.5), Inches(5.5), [
        ("Why", 0, True),
        ("・権限委譲の程度は経営判断が必要。理論設計だけでなく実プラントでの実証が必要", 0),
        ("・③-a（AI検証）とは目的が異なる。③-bは「統治・権限設計そのもの」の実証", 0),
        ("・現場主導導入の失敗リスクを、限定範囲で検証し本格展開条件を明確化", 0),
        ("・AMP構築PJのPDCA実装着手（今年度ゴール）の先行実証場として活用", 0),
    ], 10)

    add_bullet_box(slide, Inches(5.0), Inches(0.95), Inches(4.5), Inches(5.5), [
        ("What", 0, True),
        ("・定義：デジタル施策実行を題材に、権限委譲とKPI管理の運用モデルを1拠点・一定期間試行", 0),
        ("・スコープ：", 0),
        ("　- 権限委譲マトリクスの先行適用", 0),
        ("　- 発電所長主導の施策選定・実行", 0),
        ("　- KPI達成PDCA（AMP連携）", 0),
        ("　- エスカレーション・支援ラインの検証", 0),
        ("・③-aとの関係：別トラック。同一拠点も選択肢（後述Options）", 0),
    ], 10)


def slide_model_how(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "②モデルプラント設計｜How（検証設計）", slide_no=n)

    data = [
        ["#", "設計項目", "内容（案）"],
        ["1", "選定基準", "発電種別・設備規模＋発電所長の意欲・経営課題整合・DX受入体制・検証しやすさ"],
        ["2", "拠点Options", "A同一拠点（③-aと③-b）/ B別拠点 / C段階的（③-a後に③-b）→ 推奨：C（リスク分散）"],
        ["3", "権限試行", "委譲マトリクス先行適用。Level1限定→Level2標準の段階的拡大"],
        ["4", "KPI設計", "事業成果（収支・稼働率・安全）＋施策定着度。モニタリング：月次（PMO）/四半期（CJPO）"],
        ["5", "題材施策", "A.自立運営DB B.クイックWin C.収支PDCA（AMP連携）— 最低1施策、推奨2施策"],
        ["6", "プロセス", "Rev.4クイック施策ラインに「権限運用検証」トラック追加。Plan→Build→Run→Improve"],
        ["7", "成功基準", "KPI達成率≥目標80% / 施策定着率≥70% / エスカレーション適切処理率100%"],
        ["8", "撤退基準", "3ヶ月連続KPI未達→権限引き戻し検討。重大安全・コンプラ違反→即時停止"],
        ["9", "期間", "2026/10選定・体制整備 → 2026/11–2027/3実証（6ヶ月）"],
    ]
    add_table(slide, 10, 3, Inches(0.2), Inches(0.9), Inches(9.6), Inches(5.9), data)


def slide_when_who(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "②モデルプラント｜When / Who ＋ ③-aとの並行関係", slide_no=n)

    # Timeline
    tl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.95), Inches(9.4), Inches(1.8))
    tl.fill.solid()
    tl.fill.fore_color.rgb = IBM_LIGHT
    tl.line.color.rgb = IBM_BLUE
    tltf = tl.text_frame
    tltf.text = "When（スケジュール案）"
    tltf.paragraphs[0].font.bold = True
    tltf.paragraphs[0].font.size = Pt(10)

    milestones = [
        ("8/7", "v3共有"),
        ("8/12", "PO報告"),
        ("8/13", "経営層DC"),
        ("8/19", "DGD統合"),
        ("9月", "③-a開始"),
        ("10月", "③-b選定"),
        ("11月", "③-b実証"),
        ("3月", "評価・展開"),
    ]
    for i, (when, what) in enumerate(milestones):
        left = Inches(0.4 + i * 1.15)
        m = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(1.35), Inches(0.35), Inches(0.35))
        m.fill.solid()
        m.fill.fore_color.rgb = ACCENT if i >= 4 else IBM_BLUE
        m.line.fill.background()
        lbl = slide.shapes.add_textbox(left - Inches(0.15), Inches(1.75), Inches(0.7), Inches(0.7))
        ltf = lbl.text_frame
        ltf.text = f"{when}\n{what}"
        for p in ltf.paragraphs:
            p.font.size = Pt(7)
            p.alignment = PP_ALIGN.CENTER

    data = [
        ["役割", "主体", "③-a AI検証", "③-b 自立運営"],
        ["経営判断", "CJPO（CJPO担当）", "方針承認", "権限委譲範囲承認"],
        ["統合・PMO", "統括部長/長嶋/松田＋PMO拡張", "進捗・KPI管理", "権限運用ルール・接続窓口"],
        ["AMP連携", "PO（基盤部会）", "データ要件", "KPI/PDCA設計"],
        ["実証主体", "モデル発電所長", "UC検証実行", "権限運用・KPI達成"],
        ["現場実装", "G-DAC", "チェンジマネ", "定着支援"],
        ["技術支援", "DPP推進部", "開発・データ基盤", "ガードレール・カタログ"],
        ["統制", "デジタル部門", "アーキテクチャ", "セキュリティ・データ統制"],
        ["外部", "IBM", "PMO・示唆", "設計支援・ベンチマーク"],
    ]
    add_table(slide, 9, 4, Inches(0.2), Inches(2.6), Inches(9.6), Inches(4.2), data)


def slide_options(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "経営層ディスカッション向け｜判断Options", slide_no=n)

    data = [
        ["論点", "Option A", "Option B（推奨）", "Option C"],
        [
            "権限委譲の程度",
            "限定委譲\n（カタログ内のみ・小額予算）",
            "段階的委譲\n（Level1→2へ6ヶ月で拡大）",
            "拡大委譲\n（予算・選定を広範に委譲）",
        ],
        [
            "メリット",
            "リスク最小\n本社統制維持",
            "実証データ取得\nリスク/効果のバランス",
            "現場自律最大化\nスピード向上",
        ],
        [
            "リスク",
            "現場の意欲低下\n形骸化",
            "設計・運用の\n複雑性",
            "失敗・定着不全\n統制不全",
        ],
        [
            "③-b拠点",
            "③-a同一拠点",
            "段階的（③-a後）",
            "別拠点同時",
        ],
        [
            "KPI未達時",
            "即時権限引き戻し",
            "支援→改善期間→判断",
            "現場裁量維持",
        ],
        [
            "DX/AMP体制",
            "完全統合",
            "PMO共通化（推奨）",
            "完全分離",
        ],
    ]
    add_table(slide, 6, 4, Inches(0.15), Inches(0.9), Inches(9.7), Inches(4.5), data)

    add_bullet_box(slide, Inches(0.3), Inches(5.55), Inches(9.4), Inches(1.2), [
        ("IBM推奨：Option B（段階的委譲）＋ PMO共通化 ＋ ③-bは③-a開始後に段階選定。8/13経営層で委譲Level1の範囲確定を最優先判断事項とする", 0, True),
    ], 10)


def slide_next_steps(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Open Issues ／ Next Steps", slide_no=n)

    add_bullet_box(slide, Inches(0.4), Inches(0.95), Inches(4.5), Inches(5.5), [
        ("経営判断を仰ぎたい論点", 0, True),
        ("1. 権限委譲Level1の範囲（ツール選定・予算執行の上限）", 0),
        ("2. ③-bと③-aの拠点関係（同一/別/段階的）", 0),
        ("3. KPI未達時のガバナンス（介入基準・主体・頻度）", 0),
        ("4. DX推進体制とAMP構築PJの統合タイミング", 0),
        ("5. 10月組織改編との整合（体制案への影響）", 0),
    ], 10)

    add_bullet_box(slide, Inches(5.0), Inches(0.95), Inches(4.5), Inches(5.5), [
        ("Next Steps", 0, True),
        ("8/7（金）　本Rev.3を共有、簡易フィードバック受領", 0),
        ("8/12　定例：PO向け整理報告", 0),
        ("8/13〜　経営層DC：Options提示→Level1確定", 0),
        ("8/19　DGD関連資料へ統合", 0),
        ("9月　③-a選定・体制整備", 0),
        ("10月　③-b選定・実証設計確定", 0),
        ("IBM　体制図詳細化・Rev.4機能①–⑥対応表を継続", 0),
    ], 10)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs, 1)
    slide_purpose(prs, 2)
    slide_process(prs, 3)
    slide_as_is_gap(prs, 4)
    slide_why(prs, 5)
    slide_rev4_comments(prs, 6)
    slide_org_chart(prs, 7)
    slide_authority_matrix(prs, 8)
    slide_amp_connection(prs, 9)
    slide_model_why_what(prs, 10)
    slide_model_how(prs, 11)
    slide_when_who(prs, 12)
    slide_options(prs, 13)
    slide_next_steps(prs, 14)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
