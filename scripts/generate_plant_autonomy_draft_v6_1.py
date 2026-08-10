#!/usr/bin/env python3
"""Generate plant autonomy draft v6.1 — Fact→Issue→To-Be→Approach with IBM hypothesis."""

import sys
from pathlib import Path

_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from _config import output_paths

OUT_PATHS = output_paths("plant-autonomy_operating-model_draft_v6.2.pptx")

IBM_BLUE = RGBColor(0x05, 0x3F, 0x87)
IBM_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
ORANGE = RGBColor(0xC4, 0x50, 0x00)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_LINE = RGBColor(0x2E, 0x7D, 0x32)
TEXT = RGBColor(0x33, 0x33, 0x33)
FOOTER = "クライアント様 発電所自立経営 定義・IBM仮説（たたき台 v6.2）｜2026年8月7日"


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_section_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.20), Inches(12.70), Inches(0.61))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_key_message(slide, text, height=0.75):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.85), Inches(12.70), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_source_note(slide, text, top=6.75):
    box = slide.shapes.add_textbox(Inches(0.50), Inches(top), Inches(12.33), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.font.italic = True


def add_footer(slide, section, num):
    f = slide.shapes.add_textbox(Inches(0.50), Inches(7.18), Inches(10.50), Inches(0.25))
    f.text_frame.text = f"{FOOTER}　｜　{section}"
    f.text_frame.paragraphs[0].font.size = Pt(8)
    f.text_frame.paragraphs[0].font.color.rgb = GRAY
    n = slide.shapes.add_textbox(Inches(12.60), Inches(7.18), Inches(0.40), Inches(0.25))
    n.text_frame.text = str(num)
    n.text_frame.paragraphs[0].font.size = Pt(8)
    n.text_frame.paragraphs[0].font.color.rgb = GRAY
    n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, font_size=11):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, bold = item if isinstance(item, tuple) else (item, False)
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.font.bold = bold
        p.space_after = Pt(3)


def add_round_box(slide, left, top, width, height, fill, line, items, font_size=10):
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape = slide.shapes[-1]
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    add_bullets(slide, left + Inches(0.18), top + Inches(0.08), width - Inches(0.36), height - Inches(0.16), items, font_size)


def add_table(slide, data, left, top, width, height, font_size=9):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = IBM_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True


def build_title_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    t1 = slide.shapes.add_textbox(Inches(0.32), Inches(0.94), Inches(12.70), Inches(2.81))
    t1.text_frame.text = "クライアント様\n\n発電所自立経営に向けた\n「定義」と「IBM仮説」（たたき台）"
    for p in t1.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE
    t2 = slide.shapes.add_textbox(Inches(0.44), Inches(4.24), Inches(6.04), Inches(0.63))
    t2.text_frame.text = "2026年8月7日"
    t2.text_frame.paragraphs[0].font.size = Pt(14)
    t3 = slide.shapes.add_textbox(Inches(9.27), Inches(5.68), Inches(2.71), Inches(0.63))
    t3.text_frame.text = "日本アイ・ビー・エム株式会社"
    t3.text_frame.paragraphs[0].font.size = Pt(12)
    add_footer(slide, "表紙", slide_no)


def build_purpose_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "本資料の目的と読み方")
    add_key_message(
        slide,
        "クライアントの資料・ヒアリング事実を起点に、IBMとして「発電所自立経営とは何か」の仮説を提示する",
    )
    data = [
        ["読み方", "内容"],
        ["Fact", "クライアント資料・8/5ヒアリングで確認できた事実"],
        ["Issue", "事実から見える構造的課題（定義が進まない理由）"],
        ["IBM仮説", "現時点でのIBMの定義（発電所自立経営とは何か）"],
        ["To-Be", "あるべき姿（クライアント既存構想との接続）"],
        ["Approach", "仮説を実現する方策（成熟度×ガバナンス×③-b実証）"],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.75), Inches(12.33), Inches(2.35), 10)
    add_round_box(
        slide,
        Inches(0.50),
        Inches(4.30),
        Inches(12.33),
        Inches(1.05),
        RGBColor(0xFF, 0xF3, 0xE0),
        ORANGE,
        [
            ("位置づけ", True),
            "・IBMの仮説（たたき台）。エグゼクティブサマリーは本編確定後に別途作成",
            "・8/13のねらい：IBM仮説へのフィードバック ＋ 方策の優先順位合意",
        ],
        10,
    )
    add_footer(slide, "目的", slide_no)


def build_fact_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "1. Fact｜クライアントの現状（確認できた事実）")
    add_key_message(slide, "方針・責任・取組はすでに動いている。一方で「現場が何をどこまで決められるか」は未整理")
    data = [
        ["領域", "ファクト", "出典"],
        [
            "組織・責任",
            "4月より発電所長が経営責任者に。CJPO（CJPO）から数値目標を受領",
            "8/5ヒアリング 1-5",
        ],
        [
            "AMP",
            "目指す姿＝発電所自律運営でユニット価値最大化。"
            "今年度ゴール＝2027年度事業計画プロセスへのPDCA実装着手",
            "8/5ヒアリング 1-5\nAMP構築PJ資料",
        ],
        [
            "DX戦略",
            "発電所自立経営は上位戦略の一つ。DXは事業戦略実現の「手段」",
            "DX戦略策定pptx\n7/30ご案内文",
        ],
        [
            "推進体制",
            "Rev.4：発電所長＝意見提示・整合判断・先行検証。"
            "モデルPL＝AI導入・技術検証が中心",
            "Rev.4（推進体制案）",
        ],
        [
            "モデルPL",
            "③-a AI検証のみ計画。UC候補・拠点選定は未着手。"
            "「自律運営PLではない」と明言",
            "8/5ヒアリング 追加-2,3\nRev.4",
        ],
        [
            "O&ME",
            "AP-1：自律的AM実行・迅速意思決定プロセス構築が目標",
            "DPP推進部2026 AP",
        ],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.65), Inches(12.33), Inches(4.85), 8)
    add_source_note(slide, "※ 上表はIBMが8/5ヒアリングおよびクライアントご共有資料から整理した事実。解釈・仮説は次スライド以降")
    add_footer(slide, "Fact", slide_no)


def build_issue_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "2. Issue｜ファクトから見える課題")
    add_key_message(
        slide,
        "「経営責任は移ったが、経営の仕組みはまだ本社主導」のまま。"
        "だからモデルPLのAI議論だけでは自立経営は前に進まない",
    )
    data = [
        ["#", "Issue（構造的課題）", "根拠（Fact）"],
        [
            "1",
            "権限の範囲・成熟度が未定義\n（何を現場が決められるか／いつ広げるか不明）",
            "4月移行済み vs Rev.4上の発電所長役割が限定的\n初の試みで現場・本社とも不慣れ",
        ],
        [
            "2",
            "KPI管理（本社）と実行権限（現場）の非対称",
            "Rev.4機能⑥：KPI主担=発電所 vs 本社集中方針の矛盾",
        ],
        [
            "3",
            "モデルPL＝AI技術検証に限定\n自立経営の運用モデル実証場がない",
            "ヒアリング追加-3／Rev.4のモデルPL定義",
        ],
        [
            "4",
            "DX推進体制とAMP構築PJが別体系\n（合流前提だが接続設計なし）",
            "Rev.4／AMP PJ群／DX戦略が並走",
        ],
        [
            "5",
            "「困らないサポート」が体制に未制度化",
            "AMP 1-5の要請 vs Rev.4にエスカレーション未記載",
        ],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.65), Inches(12.33), Inches(4.05), 9)
    add_round_box(
        slide,
        Inches(0.50),
        Inches(5.85),
        Inches(12.33),
        Inches(0.75),
        IBM_LIGHT,
        ACCENT,
        [
            ("So What", True),
            "・Issueの本質＝「責任だけ現場、判断・投資・改善の仕組みは本社」のギャップ",
            "・権限委譲の論点は「どこから・どこまで」＋成熟度の定義・モニタリングに帰着",
        ],
        10,
    )
    add_footer(slide, "Issue", slide_no)


def build_hypothesis_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "3. IBM仮説｜発電所自立経営とは何か")
    add_key_message(
        slide,
        "発電所長がCJPO/KPI目標のもと、権限の範囲内で収支・運営・改善を自ら回す"
        "「現場CEO型」の運営モデル。権限委譲は成熟度に応じて段階的に拡大",
        0.85,
    )
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(1.55), Inches(6.00), Inches(2.75)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = IBM_BLUE
    def_box.line.color.rgb = IBM_BLUE
    def_text = slide.shapes.add_textbox(Inches(0.70), Inches(1.65), Inches(5.60), Inches(2.55))
    tf = def_text.text_frame
    tf.word_wrap = True
    lines = [
        ("IBMの定義（仮説）", True),
        "",
        "発電所自立経営とは——",
        "",
        "① CJPO/KPI目標のもと、発電所長が「社長のような役割」で",
        "　 収支・運営・改善を現場で判断・実行する",
        "",
        "② ガバナンスは成熟度に応じて設計する",
        "　 起点：本社が手綱を強めに握り、その範囲で現場裁量",
        "　 成熟：実績・定着を見て権限委譲を段階的に拡大",
        "",
        "③ DX/AIは①②を支える手段（後述：効率化→高度化）",
    ]
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, bold = item if isinstance(item, tuple) else (item, False)
        p.text = text
        p.font.size = Pt(10 if i > 0 else 11)
        p.font.bold = bold
        p.font.color.rgb = WHITE
        p.space_after = Pt(1)
    add_round_box(
        slide,
        Inches(6.70),
        Inches(1.55),
        Inches(6.13),
        Inches(2.75),
        GREEN_BG,
        GREEN_LINE,
        [
            ("クライアント既存構想との接続", True),
            "・AMP：発電所主体の収支計画PDCA（1-5）",
            "・DX戦略：自立経営領域の実現手段（DX戦略pptx）",
            "・O&ME AP-1：自律的AM・迅速意思決定（DPP AP）",
            "",
            ("≠ 以下ではない", True),
            "・一度きりの権限委譲／AI導入だけ／③-a技術PoC",
            "",
            ("＝ 以下である", True),
            "・成熟度×権限×KPI×PDCAの運用モデル",
        ],
        9,
    )
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(4.45), Inches(12.33), Inches(1.35))
    quote.fill.solid()
    quote.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    quote.line.color.rgb = ORANGE
    add_bullets(
        slide,
        Inches(0.70),
        Inches(4.55),
        Inches(11.90),
        Inches(1.15),
        [
            ("③-b 自立運営モデル発電所 ／ DX・AIの位置づけ", True),
            "・③-b：上記運用モデル（成熟度×権限×PDCA）を1拠点で試す場",
            "",
            ("DX・AIの3層（IBM仮説）", True),
            "　a. 手段：①②（自立経営の運営モデル）を実現するための道具",
            "　b. 業務効率化：作業時間・コストの削減",
            "　c. 業務高度化：暗黙知→形式知→AI化→AIと対話しながら業務",
            "　　　　　　　　（若手がベテラン知見をレバレッジ、品質向上）",
            "　※ その先の役割変革（やることが変わる）は本資料では深掘りしない",
        ],
        9,
    )
    add_source_note(slide, "根拠：8/5ヒアリング1-5、DX戦略pptx、Rev.4、DPP AP（ナレッジ形式知化）", 6.05)
    add_footer(slide, "IBM仮説", slide_no)


def build_tobe_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "4. To-Be｜あるべき姿（IBM仮説に基づく）")
    add_key_message(
        slide,
        "起点は本社ガバナンス強め＋限定裁量。成熟度が上がれば権限委譲を拡大。"
        "その設計・モニタリングがガバナンスの中核",
    )
    data = [
        ["観点", "As-Is / Gap", "To-Be（IBM仮説）"],
        [
            "経営の主体",
            "責任は現場、判断は本社寄り",
            "発電所長＝現場CEO。成熟度に応じた裁量で施策・改善を実行",
        ],
        [
            "ガバナンス\n／権限",
            "委譲範囲・拡大条件が未定義",
            "起点：本社がKPI・ガードレール・手綱を握る\n"
            "目標：成熟度評価に基づき段階的に委譲拡大",
        ],
        [
            "成熟度",
            "定義・モニタリング手段なし",
            "KPI達成・施策定着・PDCA実行力等で評価\n"
            "③-bで定義→全社展開判断の基準に",
        ],
        [
            "KPI",
            "設定（本社）と主担（Rev.4機能⑥）が混在",
            "KPI設定＝CJPO／達成責任＝発電所長／横串管理＝PMO",
        ],
        [
            "DX・AI",
            "AI検証＝技術・効率化中心",
            "手段→効率化→高度化（暗黙知の形式知化・AI対話）",
        ],
        [
            "モデルPL",
            "③-a AI検証のみ",
            "③-a（技術）＋③-b（成熟度×権限×PDCA実証）",
        ],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.65), Inches(12.33), Inches(4.05), 8)
    add_round_box(
        slide,
        Inches(0.50),
        Inches(5.85),
        Inches(12.33),
        Inches(0.85),
        IBM_LIGHT,
        ACCENT,
        [
            ("成熟度モデル（案）", True),
            "Level 0→1（起点）：本社主導＋限定委譲（カタログ内ツール・小額予算等）→ "
            "Level 2以降：実証データに基づき段階拡大（③-bで検証）",
        ],
        9,
    )
    add_footer(slide, "To-Be", slide_no)


def build_approach_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "5. Approach｜IBM仮説を実現する方策")
    add_key_message(
        slide,
        "定義合意 → ③-bで成熟度×権限×PDCAを実証 → 制度化。"
        "権限委譲は「どこから・どこまで・いつ広げるか」を③-b設計で決める",
    )
    steps = [
        ("Step 1｜定義合意", "8/7–8/13", "IBM仮説（Slide 3）への合意・修正"),
        (
            "Step 2｜③-b設計",
            "8月下旬–9月",
            "拠点選定／成熟度指標／権限の起点と目標／KPI／題材施策",
        ),
        (
            "Step 3｜③-b実証",
            "2026/10–2027/3",
            "Level 1（本社ガバナンス強め＋限定裁量）から開始。"
            "成熟度モニタリング→拡大判断",
        ),
        (
            "Step 4｜制度化",
            "2027年度〜",
            "成熟度モデル確定・Rev.4二層化・全所展開",
        ),
    ]
    y = 1.60
    for title, period, body in steps:
        add_round_box(
            slide,
            Inches(0.50),
            Inches(y),
            Inches(12.33),
            Inches(0.95),
            WHITE,
            IBM_BLUE,
            [(f"{title}（{period}）", True), body],
            9,
        )
        y += 1.05
    add_round_box(
        slide,
        Inches(0.50),
        Inches(5.85),
        Inches(12.33),
        Inches(1.05),
        RGBColor(0xFF, 0xF3, 0xE0),
        ORANGE,
        [
            ("Step 2で詰める論点（ガバナンス設計）", True),
            "① 権限委譲：どこからスタートし、どこまでを目指すか（起点＝Level 1）",
            "② 本社ガバナンス：手綱の握り方・例外ルール・エスカレーション",
            "③ 成熟度：定義（KPI・定着・PDCA実行力等）とモニタリング会議体",
            "④ 期間：技術検証（短期）vs 成熟度評価（最低1経営サイクル）",
        ],
        9,
    )
    add_footer(slide, "Approach", slide_no)


def build_positioning_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "6. 3つの取組の位置付け（IBM仮説との関係）")
    add_key_message(slide, "DX・AMP・モデルPLは別PJではなく、IBM仮説（運営モデル変革）の異なるレイヤー")
    data = [
        ["取組", "レイヤー", "IBM仮説との関係"],
        ["DX戦略", "変革テーマ", "自立経営の手段＋効率化＋高度化（3層）"],
        ["AMP構築PJ", "経営プロセス", "発電所主体PDCA＝自立経営の中核プロセス"],
        ["③-a AI検証", "技術検証", "③-bの入力。単独では自立経営にならない"],
        ["③-b 自立運営", "運営モデル実証", "IBM仮説そのものを1拠点で試す"],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.65), Inches(12.33), Inches(2.15), 10)
    add_round_box(
        slide,
        Inches(0.50),
        Inches(4.05),
        Inches(5.90),
        Inches(2.35),
        IBM_LIGHT,
        IBM_BLUE,
        [
            ("③-a AI検証", True),
            "目的：技術・導入可否（PoC）",
            "時期：2026/9〜2027/3（DX戦略）",
            "Rev.4既存トラック",
        ],
        9,
    )
    add_round_box(
        slide,
        Inches(6.70),
        Inches(4.05),
        Inches(6.13),
        Inches(2.35),
        GREEN_BG,
        GREEN_LINE,
        [
            ("③-b 自立運営モデル実証", True),
            "目的：成熟度×権限×PDCAの運用モデル",
            "起点：本社ガバナンス強め→段階拡大",
            "IBM新提案。AMP PDCA先行実証",
        ],
        9,
    )
    add_footer(slide, "位置付け", slide_no)


def build_prove_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "7. ③-bで何を実証するか")
    add_key_message(
        slide,
        "経営の仕組み（成熟度×権限×PDCA）が先。"
        "AIは手段→効率化→高度化の順で価値が立つ",
    )
    data = [
        ["順", "実証項目", "IBM仮説上の意味"],
        ["1", "ガバナンス設計（どこから・どこまで）", "Issue#1：権限＋成熟度の定義"],
        ["2", "成熟度モニタリング（評価・拡大判断）", "段階的委譲の根拠データ"],
        ["3", "収支PDCA（AMP接続）", "AMP今年度ゴールの先行実証"],
        ["4", "KPI達成と本社統制の両立", "Issue#2：非対称の解消"],
        [
            "5",
            "DX・AI（3層）",
            "a.手段 b.効率化 c.高度化\n"
            "（暗黙知→形式知→AI対話→品質向上）",
        ],
    ]
    add_table(slide, data, Inches(0.50), Inches(1.65), Inches(12.33), Inches(2.85), 9)
    add_round_box(
        slide,
        Inches(0.50),
        Inches(4.70),
        Inches(12.33),
        Inches(1.35),
        IBM_LIGHT,
        ACCENT,
        [
            ("成功基準（案）", True),
            "・成熟度指標：KPI達成率・施策定着率・PDCA実行力・権限設計の妥当性",
            "・最低1経営サイクル（6〜12ヶ月）で評価（技術PoCの短期サイクルとは別）",
            "・AI題材：効率化（Quick Win等）＋高度化（ナレッジ形式知化・DPP AP）",
            "・③-a成果は③-bの「手段・高度化」層に接続",
        ],
        9,
    )
    add_footer(slide, "実証内容", slide_no)


def build_next_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_section_title(slide, "8. IBMの主張と次のステップ")
    add_key_message(
        slide,
        "発電所自立経営＝一度きりの権限委譲でもAI導入でもなく、"
        "成熟度に応じて進化する現場CEO型の運営モデル",
    )
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(1.55), Inches(12.33), Inches(0.95))
    quote.fill.solid()
    quote.fill.fore_color.rgb = IBM_BLUE
    quote.line.color.rgb = IBM_BLUE
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.68), Inches(11.80), Inches(0.72))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = (
        "IBMの主張：クライアントは「現場に経営責任を移す」方向に進んでいる。"
        "次に必要なのは、起点（本社ガバナンス強め）と目標（成熟度に応じた委譲拡大）を定義し、"
        "③-bで実証・モニタリングすること。"
        "AIは手段に加え、効率化・高度化（暗黙知の形式知化）として自立経営を支える。"
    )
    for p in tf.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    data = [
        ["時期", "アクション", "ねらい"],
        ["8/7–8/12", "PO・プロジェクトオーナーですり合わせ", "Fact/Issue/IBM仮説へのフィードバック"],
        ["8/13", "経営層ディスカッション", "IBM仮説合意 ＋ Approach優先順位"],
        ["9月以降", "③-b具体設計", "権限/KPI/拠点/題材/体制案"],
        ["来週", "エグゼクティブサマリー", "本編確定後に別途"],
    ]
    add_table(slide, data, Inches(0.50), Inches(2.60), Inches(12.33), Inches(2.35), 10)
    add_source_note(slide, "参照：8/5ヒアリングシート、Rev.4、DX戦略策定pptx、7/30ご案内、AMP構築PJ、DPP 2026 AP", 5.20)
    add_footer(slide, "Next", slide_no)


def build_eof_slide(prs, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.20), Inches(12.70), Inches(0.61))
    box.text_frame.text = "End of Document"
    box.text_frame.paragraphs[0].font.size = Pt(24)
    box.text_frame.paragraphs[0].font.bold = True
    add_footer(slide, "End", slide_no)


def main():
    dst = Presentation()
    dst.slide_width = Inches(13.33)
    dst.slide_height = Inches(7.50)

    build_title_slide(dst, 1)
    build_purpose_slide(dst, 2)
    build_fact_slide(dst, 3)
    build_issue_slide(dst, 4)
    build_hypothesis_slide(dst, 5)
    build_tobe_slide(dst, 6)
    build_approach_slide(dst, 7)
    build_positioning_slide(dst, 8)
    build_prove_slide(dst, 9)
    build_next_slide(dst, 10)
    build_eof_slide(dst, 11)

    for out in OUT_PATHS:
        out.parent.mkdir(parents=True, exist_ok=True)
        dst.save(str(out))
        print(f"Saved: {out}")
    print(f"Total slides: {len(dst.slides)}")


if __name__ == "__main__":
    main()
