#!/usr/bin/env python3
"""Generate JERA SCN deck anchored on EBITDA tree (knowledge/patterns/jera-scn-ebitda-tree.md)."""

import sys
from pathlib import Path

_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

from _config import output_paths

OUT_PATHS = output_paths("JERA_SCN_ValueCreationTree_v2.3.pptx")

IBM_BLUE = RGBColor(0x05, 0x3F, 0x87)
IBM_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xC4, 0x50, 0x00)
ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_LINE = RGBColor(0x2E, 0x7D, 0x32)
CAP_FILL = RGBColor(0xDE, 0xEB, 0xF7)
EN_FILL = RGBColor(0xF5, 0xF5, 0xF5)
DECISION_FILL = RGBColor(0xFF, 0xF8, 0xE1)
VD_FILL = RGBColor(0xE0, 0xF2, 0xF1)
FOOTER = "JERA｜SCN（Value Creation Tree）v2.3｜2026年8月13日"
FONT_BODY = 14
FONT_TITLE = 22
FONT_SUB = 12
FONT_FOOTER = 9
FONT_BOX = 12

# 有価証券報告書（第10期・2025年3月期）個別財務諸表 P.140–142（印字）
# 会計基準：JGAAP・親会社単体。連結IFRSセグメント（P.85）とは混在しない。
YUHO_FY25 = {
    "他社販売電力料": "3,752,519",
    "電気事業雑収益": "88,871",
    "燃料費": "2,887,113",
    "修繕費": "113,056",
    "消耗品費": "16,182",
    "委託費": "30,397",
    "賃借料": "5,091",
    "給料手当": "30,860",
    "厚生費": "5,932",
    "他社購入電力料": "235,909",
    "減価償却費_個別": "116,753",
    "減価償却費_連結参考": "199,593",
    "汽力発電費合計": "3,292,729",
    "電気事業営業利益_個別": "124,483",
    "セグメント利益_連結参考": "124,324",
    "一般管理費": "55,559",
    "接続供給託送料": "91,159",
    "事業税": "38,778",
    "その他未掲載合計": "185,496",
    "ガス供給収益": "421,320",
    "ガス供給費用": "415,958",
}
YUHO_SOURCE = "securities_report2506.pdf P.140–142（個別・印字）"
YUHO_SCOPE_NOTE = (
    "個別JGAAP・電気事業のみ｜ガス供給事業・一般管理費等185,496Mは対象外（抜粋）"
)

# JPJ DGD v1.0 / プログラム定義書（20260813）— Value-1 ベースラインとは別層（増分 KGI）
JPJ_KGI = {
    "焚替損失回避_2027": "86.51億円/年",
    "焚替損失回避_2030": "109億円/年",
    "運用費削減": "5億円/年",
    "工数_2027": "84.1名＋390名",
    "工数_2030": "569名＋390名",
    "内訳_保全": "79.5億",
    "内訳_運転": "29.3億",
    "内訳_資材SC": "0.37億",
    "GDAC実績_2025": "18億円",
}
JPJ_SOURCE = "JPJ DGD v1.0 §7–8／20260813_プログラム定義書"


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_section_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(0.18), Inches(12.70), Inches(0.55))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(FONT_TITLE)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_subtitle(slide, text, top=0.72):
    box = slide.shapes.add_textbox(Inches(0.32), Inches(top), Inches(12.70), Inches(0.40))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(FONT_SUB)
    p.font.color.rgb = GRAY


def add_footer(slide, section, num):
    f = slide.shapes.add_textbox(Inches(0.40), Inches(7.15), Inches(11.80), Inches(0.25))
    f.text_frame.text = f"{FOOTER}　｜　{section}"
    f.text_frame.paragraphs[0].font.size = Pt(FONT_FOOTER)
    f.text_frame.paragraphs[0].font.color.rgb = GRAY
    n = slide.shapes.add_textbox(Inches(12.55), Inches(7.15), Inches(0.45), Inches(0.25))
    n.text_frame.text = str(num)
    n.text_frame.paragraphs[0].font.size = Pt(FONT_FOOTER)
    n.text_frame.paragraphs[0].font.color.rgb = GRAY
    n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def style_cell(cell, font_size, header=False, center=False):
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(3)
    cell.margin_top = cell.margin_bottom = Pt(2)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE if header else TEXT
        p.font.bold = header
        if center:
            p.alignment = PP_ALIGN.CENTER
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = IBM_BLUE


def add_table(slide, data, left, top, width, height, font_size=FONT_BODY, col_widths=None, center_cols=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    center_cols = center_cols or set()
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            style_cell(cell, font_size, header=(r == 0), center=(c in center_cols))
    return table


def new_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
    return slide


def add_box(slide, left, top, width, height, text, fill, line, font_size=FONT_BOX, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_oval(slide, left, top, width, height, text, fill, line, font_size=FONT_BOX, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(2 if bold else 1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = TEXT if fill != IBM_BLUE else WHITE
    p.alignment = PP_ALIGN.CENTER
    return shape


def connect(slide, x1, y1, x2, y2, color=GRAY):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.25)
    return line


def build_title(prs, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    t = slide.shapes.add_textbox(Inches(0.50), Inches(1.20), Inches(12.33), Inches(2.00))
    t.text_frame.text = "JERA Strategic Capability Network\nValue Creation Tree SCN（v2.3）"
    for p in t.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE
    s = slide.shapes.add_textbox(Inches(0.50), Inches(3.40), Inches(12.00), Inches(1.70))
    s.text_frame.word_wrap = True
    s.text_frame.text = (
        "6層：EBITDA → Value Driver → Decision → Process → Capability → Digital/AI\n"
        "Outcome＝有報124,483M（Baseline）｜KGI＝焚替回避86.51/109億（増分・Forward）\n"
        "knowledge/patterns/jera-scn-ebitda-tree.md v1.3"
    )
    s.text_frame.paragraphs[0].font.size = Pt(FONT_BODY)
    s.text_frame.paragraphs[0].font.color.rgb = GRAY
    add_footer(slide, "表紙", n)


def build_rationale(prs, n):
    slide = new_slide(
        prs,
        "SCN再設計のねらい",
        "Value Creation Tree：EBITDA→Value Driver→Decision→Capability→Enabler（BPR先行）",
    )
    add_table(
        slide,
        [
            ["層", "役割", "有報・Operating Model接続"],
            ["Value", "EBITDA最大化（Outcome）\nDX/AIはValueに置かない", "Baseline：124,483M＋D&A\nKGI（増分）：86.51→109億"],
            ["Value Driver", "有報費目の因数分解\n（熱効率・送電量・LCC等）", "燃料費→熱効率→運転最適化"],
            ["Decision", "誰が何を判断するか\n（所長の経営判断）", "運転計画見直し・修繕優先・AI採否"],
            ["Capability", "Business / Enterprise\n「～できる能力」", "Monitor KPI＋Decision KPI"],
            ["Enabler", "BPR→Process→Technology\n（DPP教訓：AI先行×）", "AMP/DGD/BPR/原価/データ基盤/資材DX"],
        ],
        Inches(0.30),
        Inches(1.05),
        Inches(12.70),
        Inches(3.15),
        font_size=11,
        col_widths=[Inches(1.25), Inches(5.45), Inches(6.00)],
    )
    note = add_box(
        slide,
        Inches(0.30),
        Inches(4.35),
        Inches(12.70),
        Inches(0.95),
        "報告4行：①Value Driver＝有報のどの枝か　②Decision＝所長/本社の判断　"
        "③Capability＝Monitor KPI　④Enabler＝BPR→施策（判断→行動→KPI→EBITDA）",
        ORANGE_BG,
        ORANGE,
        FONT_BODY,
        True,
    )
    add_footer(slide, "ねらい", n)


def build_yuho_mapping(prs, n):
    slide = new_slide(
        prs,
        "有報費目との対応（第10期・2025年3月期）",
        f"出典：{YUHO_SOURCE}（単位：百万円）｜{YUHO_SCOPE_NOTE}",
    )
    y = YUHO_FY25
    add_table(
        slide,
        [
            ["McK/SCN", "有報費目", "符号", "FY25", "Value Driver", "Cap"],
            ["送電×スプレッド", "他社販売電力料", "＋", y["他社販売電力料"], "送電MWh×スプレッド", "C-A1"],
            ["付加収益", "電気事業雑収益", "＋", y["電気事業雑収益"], "付加収益機会", "Cap-A"],
            ["O&M・変動費", "燃料費", "－", y["燃料費"], "熱効率・運転パターン", "C-A2"],
            ["メンテ", "修繕費＋消耗品費", "－", f'{y["修繕費"]}＋{y["消耗品費"]}', "LCC・設備リスク", "C-B1"],
            ["委託・賃借", "委託費＋賃借料", "－", f'{y["委託費"]}＋{y["賃借料"]}', "生産性・FTE", "C-B3"],
            ["人件費", "給料＋厚生", "－", f'{y["給料手当"]}＋{y["厚生費"]}', "標準化・定着", "C-B2/X2*"],
            ["他社電源", "他社購入電力料", "－", y["他社購入電力料"], "需給・調達", "Cap-A"],
            ["D&A加算", "減価償却費（個別）", "＋", y["減価償却費_個別"], "—", "—"],
            ["Value-1", "電気事業営業利益", "＝", y["電気事業営業利益_個別"], "Outcome", "Value-1"],
        ],
        Inches(0.12),
        Inches(1.05),
        Inches(13.05),
        Inches(5.35),
        font_size=10,
        col_widths=[Inches(1.5), Inches(3.2), Inches(0.4), Inches(1.5), Inches(2.2), Inches(0.75)],
        center_cols={2},
    )
    add_box(
        slide,
        Inches(0.18),
        Inches(6.55),
        Inches(12.95),
        Inches(0.55),
        f"未掲載185,496M・ガス事業対象外｜費目→Value Driver→Capability（直結しない）",
        EN_FILL,
        GRAY,
        9,
    )
    add_footer(slide, "有報対応", n)


def build_ebitda_tree(prs, n):
    slide = new_slide(
        prs,
        "Value Creation Tree（有報 × Value Driver）",
        f"燃料費例：燃料費→熱効率→運転最適化→Decision→C-A2｜{YUHO_SCOPE_NOTE}",
    )
    y = YUHO_FY25
    items = [
        (f"Value-1：EBITDA\n{y['電気事業営業利益_個別']}M（個別）", "Outcome", IBM_BLUE, WHITE, 0.55, 0.98, 12.0, 0.52),
        (f"＋ 収益 {y['他社販売電力料']}M＋{y['電気事業雑収益']}M", "有報＋", IBM_LIGHT, IBM_BLUE, 0.55, 1.58, 12.0, 0.42),
        (f"－ 燃料費 {y['燃料費']}M", "有報－", CAP_FILL, IBM_BLUE, 0.55, 2.08, 2.5, 0.55),
        (f"Value Driver\n熱効率・運転パターン", "VD", VD_FILL, GREEN_LINE, 3.15, 2.08, 2.5, 0.55),
        (f"Decision\n運転計画見直し", "所長", DECISION_FILL, ORANGE, 5.75, 2.08, 2.5, 0.55),
        (f"Capability\nC-A2 運転最適化", "Business", CAP_FILL, IBM_BLUE, 8.35, 2.08, 2.5, 0.55),
        (f"Enabler\nBPR→AI最適化", "手段", EN_FILL, GRAY, 10.95, 2.08, 2.1, 0.55),
        (f"－ 修繕 {y['修繕費']}M＋消耗 {y['消耗品費']}M", "有報－", CAP_FILL, IBM_BLUE, 0.55, 2.78, 2.5, 0.55),
        (f"VD: LCC・リスク", "VD", VD_FILL, GREEN_LINE, 3.15, 2.78, 2.5, 0.55),
        (f"Decision: 修繕優先", "所長", DECISION_FILL, ORANGE, 5.75, 2.78, 2.5, 0.55),
        (f"C-B1 保全最適化", "Business", CAP_FILL, IBM_BLUE, 8.35, 2.78, 2.5, 0.55),
        (f"－ 人件費 / 他社購入 / 計停・計外", "その他枝", CAP_FILL, IBM_BLUE, 0.55, 3.48, 7.3, 0.55),
        (f"＋ D&A {y['減価償却費_個別']}M", "EBITDA加算", EN_FILL, GRAY, 8.05, 3.48, 4.0, 0.55),
    ]
    for text, sub, fill, line, l, t, w, h in items:
        add_box(slide, Inches(l), Inches(t), Inches(w), Inches(h), f"{text}\n{sub}" if sub else text, fill, line, 9, fill == IBM_BLUE)
    add_box(
        slide,
        Inches(0.55),
        Inches(4.25),
        Inches(12.0),
        Inches(0.45),
        "凡例：有報（結果）→ Value Driver（因数）→ Decision（判断）→ Capability → Enabler（BPR→AI）",
        ORANGE_BG,
        ORANGE,
        9,
    )
    add_footer(slide, "Value Tree", n)


def build_scn_map(prs, n):
    slide = new_slide(prs, "SCN全体像", "Value → Decision → Capability（C-0主経路）→ BPR/Enabler")
    add_oval(slide, Inches(4.60), Inches(0.88), Inches(4.20), Inches(0.62), "Value-1\nEBITDA最大化", IBM_BLUE, IBM_BLUE, 10)
    add_oval(slide, Inches(4.20), Inches(1.62), Inches(5.00), Inches(0.52), "Decision\n所長・本社の経営判断", DECISION_FILL, ORANGE, 9, False)
    add_box(
        slide,
        Inches(1.00),
        Inches(2.30),
        Inches(11.33),
        Inches(0.78),
        "【Enterprise・主経路】C-0：収支PDCA → 経営判断 → Business Capability\n"
        "Decision KPI：計画修正・優先順位｜Monitor：PDCA完遂・計画精度",
        ORANGE_BG,
        ORANGE,
        9,
        True,
    )
    caps = [
        ("Cap-A\n収益", 0.55, 3.35),
        ("Cap-B\nO&M", 3.55, 3.35),
        ("Cap-C\n停止", 6.55, 3.35),
        ("Cap-X\nEnterprise", 9.55, 3.35),
    ]
    for label, l, t in caps:
        add_oval(slide, Inches(l), Inches(t), Inches(2.85), Inches(0.88), label, CAP_FILL, IBM_BLUE, 9)
    enablers = [
        ("BPR/業務棚卸", 0.35),
        ("AMP構築PJ", 2.55),
        ("DGD/Team DX", 4.75),
        ("データ基盤", 6.95),
        ("原価管理", 9.15),
        ("資材DX", 11.35),
    ]
    for label, l in enablers:
        add_box(slide, Inches(l), Inches(4.55), Inches(1.85), Inches(0.50), label, EN_FILL, GRAY, 8)
    cx = Inches(6.70)
    connect(slide, cx, Inches(1.50), cx, Inches(2.30))
    connect(slide, cx, Inches(3.08), cx, Inches(3.35))
    for l in [1.97, 4.97, 7.97, 10.97]:
        connect(slide, cx, Inches(3.08), Inches(l), Inches(3.35))
        connect(slide, Inches(l), Inches(4.23), Inches(l + 0.3), Inches(4.55))
    add_box(
        slide,
        Inches(0.25),
        Inches(5.20),
        Inches(12.75),
        Inches(0.38),
        "【PG-1】AMP×DGD 基盤重複 → C-X1/C-0 共有ノード＋RACI分担｜自立経営＝C-0/C-X3（プログラムゴール）",
        ORANGE_BG,
        ORANGE,
        9,
        True,
    )
    add_footer(slide, "SCN全体像", n)


def build_capability_table(prs, n):
    slide = new_slide(
        prs,
        "Capability一覧（Business / Enterprise × Decision）",
        "種別：Biz=業務能力 Ent=組織能力｜○＝主担当",
    )
    headers = ["ID", "種", "Capability", "Value Driver", "Decision", "Owner", "Monitor KPI"]
    rows = [
        ["C-0", "Ent", "収支PDCA【太線】", "計画精度", "計画修正承認", "所長+本社", "PDCA完遂"],
        ["C-A1", "Biz", "需給・運転計画実行", "送電MWh", "稼働計画見直し", "所長", "稼働率"],
        ["C-A2", "Biz", "熱効率・収益確保", "熱効率", "運転計画見直し", "所長", "熱効率"],
        ["C-B1", "Biz", "保全・修繕最適化", "LCC", "修繕優先/延期", "所長", "点検INT"],
        ["C-B2", "Biz", "O&Mコスト管理", "起動費", "起動回数調整", "所長", "起動回数"],
        ["C-B3", "Biz", "委託生産性", "FTE", "委託範囲見直し", "所長", "修繕期間"],
        ["C-C1", "Biz", "計画停止最適化", "計停MWh", "定検時期", "所長", "計停日数"],
        ["C-C2", "Biz", "計画外停止防止", "計外MWh", "予兆対応優先", "所長", "計外率"],
        ["C-C3", "Biz", "早期復旧", "停止損失", "初動判断", "所長", "復旧日数"],
        ["C-X1", "Ent", "timely把握", "Input", "データ投資", "所長+本社", "整備率"],
        ["C-X2", "Ent", "標準プロセス再現", "BPR", "標準採否", "所長+CoE", "定着率"],
        ["C-X3", "Ent", "裁量実行", "権限", "ツール導入", "所長", "委譲度"],
    ]
    add_table(
        slide,
        [headers] + rows,
        Inches(0.12),
        Inches(1.05),
        Inches(13.05),
        Inches(5.85),
        font_size=9,
        col_widths=[
            Inches(0.48),
            Inches(0.38),
            Inches(2.85),
            Inches(1.15),
            Inches(1.35),
            Inches(0.75),
            Inches(1.15),
        ],
        center_cols={1},
    )
    add_footer(slide, "Capability", n)


def build_enabler_table(prs, n):
    slide = new_slide(prs, "Enabler束・BPRマッピング", "モデルプラントAI → BPR → Capability → EBITDA（施策単体×）")
    add_box(
        slide,
        Inches(0.25),
        Inches(1.02),
        Inches(12.75),
        Inches(0.42),
        "DPP教訓：AI/Technology先行× → BPR（Process/Knowledge）→ Technology の順で接続",
        ORANGE_BG,
        ORANGE,
        10,
        True,
    )
    add_table(
        slide,
        [
            ["Initiative束", "接続Cap", "BPR/KOPT", "主体"],
            ["AMP構築PJ", "C-0, C-C", "P:収支PDCA / O:責任分界", "高橋PO・長嶋PM"],
            ["発電所自立経営", "C-0, C-X3", "O:権限委譲 / P:所長PDCA", "運営統括（プログラムゴール）"],
            ["Team DX / DGD", "C-A–C, C-X", "P:業務成熟度 / T:3+", "森崎PO・手川PM"],
            ["プラントデータ基盤", "C-X1", "P:データフロー / T:構造化", "長嶋PO・市場DO"],
            ["原価管理PJ", "C-0", "P:両輪PDCA / T:S4", "酒入/森崎/行徳"],
            ["資材DX", "C-B3", "T:調達PF / P:SC業務", "鈴木PO（構想）"],
            ["設備診断支援", "C-C1, C-B1", "P:懸案一元化 / K:INT", "DPP連携"],
        ],
        Inches(0.25),
        Inches(1.55),
        Inches(12.75),
        Inches(2.85),
        font_size=10,
        col_widths=[Inches(2.0), Inches(1.6), Inches(5.8), Inches(3.35)],
    )
    add_box(
        slide,
        Inches(0.25),
        Inches(4.52),
        Inches(12.75),
        Inches(0.32),
        "【PG-1】AMP×DGD：C-X1/C-0 を共有。データ基盤とDGD施策の重複はRACIで解消",
        ORANGE_BG,
        ORANGE,
        9,
        True,
    )
    add_table(
        slide,
        [
            ["施策/BPR", "Cap", "KGI/連鎖", "所長", "本社", "DX"],
            ["業務棚卸(千葉)", "C-X2", "BPR→標準化", "○", "○", "△"],
            ["モデルプラントAI", "C-C2/B1", "BPR→AI→Cap", "○", "○", "○"],
            ["AMP全体構築部会", "C-0", "PDCA To-Be", "○", "○", "△"],
            ["ユニット別収支DB", "C-0", "Decision支援", "○", "○", "○"],
            ["計画外抑制部会", "C-C2", "予兆→109億の主因", "○", "○", "△"],
            ["設備保守DB部会", "C-C2,X1", "データ→判断", "○", "○", "○"],
            ["設備診断支援", "C-C1,B1", "INT最適化", "○", "○", "○"],
            ["J-AIME/G-DAC", "C-C2,A2", "18億実績→KGI", "△", "○", "○"],
            ["10/1組織変更", "C-X3,C-0", "権限→Decision", "○", "○", "—"],
        ],
        Inches(0.25),
        Inches(4.95),
        Inches(12.75),
        Inches(2.35),
        col_widths=[Inches(2.4), Inches(1.1), Inches(1.8), Inches(0.5), Inches(0.5), Inches(0.5)],
        center_cols={3, 4, 5},
    )
    add_footer(slide, "Enabler", n)


def build_kpi_reconciliation(prs, n):
    k = JPJ_KGI
    y = YUHO_FY25
    slide = new_slide(
        prs,
        "KPI階層：有報Baseline vs プログラムKGI",
        f"出典：有報＝{YUHO_SOURCE}｜KGI＝{JPJ_SOURCE}",
    )
    add_table(
        slide,
        [
            ["層", "指標", "数値", "SCN位置", "プログラム定義書との関係"],
            ["A Baseline", "電気事業営業利益", f"{y['電気事業営業利益_個別']}M（≈1,245億）", "Value-1 アンカー", "FY25有報＝現状Outcome"],
            ["A′ EBITDA", "営業利益＋D&A", f"{int(y['電気事業営業利益_個別'].replace(',',''))+int(y['減価償却費_個別'].replace(',','')):,}M", "Value-1補足", "EBITDA説明用（個別）"],
            ["B KGI 2027", "焚替損失回避", k["焚替損失回避_2027"], "VD→C-C/C-B増分", "≠124,483M。増分効果"],
            ["B KGI 2030", "焚替損失回避", k["焚替損失回避_2030"], "VD→C-C/C-B増分", "DGD §7 2030目標"],
            ["B KGI", "運用費削減", k["運用費削減"], "C-B2/Enabler", "AIツール等"],
            ["B 内訳", "保全/運転/資材", f"{k['内訳_保全']}/{k['内訳_運転']}/{k['内訳_資材SC']}", "C-B/C-C/C-A2", "109億の因数分解"],
            ["C 実績", "GDAC予兆監視", k["GDAC実績_2025"], "C-C2実証", "KGI達成の先行事例"],
        ],
        Inches(0.18),
        Inches(1.05),
        Inches(12.95),
        Inches(3.55),
        font_size=10,
        col_widths=[Inches(1.05), Inches(1.55), Inches(1.85), Inches(1.65), Inches(3.15)],
    )
    add_box(
        slide,
        Inches(0.18),
        Inches(4.75),
        Inches(12.95),
        Inches(1.05),
        "整合ルール：124,483M（有報）と86.51/109億（KGI）は同じ数値ではない。\n"
        "有報＝Where we are（Baseline Outcome）｜KGI＝What program delivers（Forward増分）\n"
        "McKストレッチ＝A＋Bの将来像。施策報告は「判断→行動→VD枝→KGI/有報枝」の連鎖で説明",
        ORANGE_BG,
        ORANGE,
        10,
        True,
    )
    add_table(
        slide,
        [
            ["KGI枝", "主Capability", "主Enabler", "2027→2030"],
            ["焚替・計外停止", "C-C2, C-C3", "AMP計外抑制・J-AIME", "86.51→109億"],
            ["保全・計停", "C-B1, C-C1", "設備診断・設備保守DB", "79.5億内"],
            ["運転・燃料", "C-A2", "G-DAC・運転データ", "29.3億内"],
            ["収支PDCA", "C-0", "AMP・原価管理", "全枝統合"],
        ],
        Inches(0.18),
        Inches(5.95),
        Inches(12.95),
        Inches(1.05),
        font_size=9,
        col_widths=[Inches(1.5), Inches(1.5), Inches(3.5), Inches(1.5)],
    )
    add_footer(slide, "KPI階層", n)


def build_kpi_gap(prs, n):
    slide = new_slide(prs, "KPI配置とGap対応", "Outcome / Monitor / Decision / 変革KPI｜G7＝判断→行動→KPI→EBITDA")
    add_table(
        slide,
        [
            ["レイヤ", "ノード", "指標種別", "例"],
            ["Value-1", "EBITDA", "Baseline", "124,483M＋D&A116,753M"],
            ["Value-1′", "Program KGI", "Forward増分", "焚替86.51→109億/運用費5億"],
            ["Value Driver", "熱効率等", "Driver KPI", "燃料費の因数分解"],
            ["Decision", "所長判断", "Decision KPI", "計画修正・修繕優先"],
            ["C-0", "収支PDCA", "Monitor", "PDCA完遂"],
            ["Enabler", "BPR/DX", "変革KPI", "定着・導入（Outcome×）"],
        ],
        Inches(0.30),
        Inches(1.05),
        Inches(6.20),
        Inches(2.75),
    )
    add_table(
        slide,
        [
            ["Gap", "症状", "EBITDA SCNでの修正"],
            ["G1", "組織先行・中身後", "10/1権限とC-0 To-Beを同一図でExCom合意"],
            ["G2", "PDCA3系統分裂", "C-0にAMP/SAP/Dataikuを束ねOutcome一本化"],
            ["G3", "DGD体制未起動", "Team DX PMOをSCN右肩に明示"],
            ["G4", "地点未決", "C-Cに1拠点を太線"],
            ["G5", "KPI混同", "Baseline/KGIを3層分離"],
            ["G6", "オーナー分散", "BPR→C-X2に集約"],
            ["G7", "DPP再演", "判断→行動→KPI→EBITDA枝＋BPR先行"],
            ["G8", "PG-1重複", "AMP×DGD→C-X1/C-0共有+RACI"],
            ["V-02", "自立未定義", "C-X3とC-0 To-Be同一図"],
        ],
        Inches(6.70),
        Inches(1.05),
        Inches(6.30),
        Inches(3.55),
    )
    add_footer(slide, "KPI・Gap", n)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 1
    build_title(prs, n)
    n += 1
    build_rationale(prs, n)
    n += 1
    build_yuho_mapping(prs, n)
    n += 1
    build_ebitda_tree(prs, n)
    n += 1
    build_scn_map(prs, n)
    n += 1
    build_capability_table(prs, n)
    n += 1
    build_enabler_table(prs, n)
    n += 1
    build_kpi_reconciliation(prs, n)
    n += 1
    build_kpi_gap(prs, n)
    for path in OUT_PATHS:
        prs.save(str(path))
        print(f"Wrote {path}")


if __name__ == "__main__":
    main()
