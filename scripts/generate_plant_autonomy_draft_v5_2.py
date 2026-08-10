#!/usr/bin/env python3
"""Generate plant autonomy draft v5.2 by cloning v5.1 slides in order + 2 new slides."""

import sys
from pathlib import Path

_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

import shutil
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from _config import OUTPUT_DIR, output_paths

SRC = OUTPUT_DIR / "plant-autonomy_operating-model_draft_v5.1.pptx"
OUT_PATHS = output_paths("plant-autonomy_operating-model_draft_v5.2.pptx")

IBM_BLUE = RGBColor(0x05, 0x3F, 0x87)
IBM_LIGHT = RGBColor(0xE8, 0xF0, 0xFA)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x78, 0xD4)
ORANGE = RGBColor(0xC4, 0x50, 0x00)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_LINE = RGBColor(0x2E, 0x7D, 0x32)
TEXT = RGBColor(0x33, 0x33, 0x33)
FOOTER = "クライアント様 発電所自立経営 体制・モデルプラント設計 骨子（たたき台 v5.2）｜2026年8月5日時点"


def clone_slide(src_prs, src_idx, dst_prs):
    src = src_prs.slides[src_idx]
    layout = dst_prs.slide_layouts[0]
    dst = dst_prs.slides.add_slide(layout)
    for shape in src.shapes:
        el = deepcopy(shape.element)
        dst.shapes._spTree.insert_element_before(el, "p:extLst")
    if src.background.fill.type:
        dst.background.fill.solid()
        dst.background.fill.fore_color.rgb = src.background.fill.fore_color.rgb
    return dst


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_label(slide, text):
    box = slide.shapes.add_textbox(Inches(0.50), Inches(0.30), Inches(11.50), Inches(0.28))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT
    p.font.bold = True


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.50), Inches(0.55), Inches(12.30), Inches(0.55))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE


def add_footer(slide, section, num):
    f = slide.shapes.add_textbox(Inches(0.50), Inches(7.18), Inches(10.50), Inches(0.25))
    f.text_frame.text = f"{FOOTER}　｜　{section}"
    f.text_frame.paragraphs[0].font.size = Pt(8)
    f.text_frame.paragraphs[0].font.color.rgb = GRAY
    n = slide.shapes.add_textbox(Inches(12.60), Inches(7.18), Inches(0.40), Inches(0.25))
    n.text_frame.text = str(num)
    n.text_frame.paragraphs[0].font.size = Pt(8)
    n.text_frame.paragraphs[0].font.color.rgb = GRAY
    n.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, font_size=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, bold = item if isinstance(item, tuple) else (item, False)
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.font.bold = bold
        p.space_after = Pt(3)


def add_table(slide, data, left, top, width, height):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = IBM_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True


def build_maturity_slide(slide, slide_no):
    set_bg(slide)
    add_label(slide, "第1部｜エグゼクティブ")
    add_title(slide, "権限成熟度モデル（案）｜Level 0〜3")
    data = [
        ["Level", "名称", "委譲範囲（発電所長）", "本社/本社が保持", "トラックB実証での位置づけ"],
        ["0", "本社主導", "ツール選定・予算・KPI設定すべて本社", "—", "現状の一部（移行前）"],
        ["1", "限定委譲", "カタログ内ツール・小額予算・所内運用設計", "KPI設定・基準・大型投資", "★ 実証開始点（IBM推奨）"],
        ["2", "標準委譲", "Level1＋一定額予算執行・施策選定", "KPI設定・セキュリティ統制", "6ヶ月後の拡大検討"],
        ["3", "拡大委譲", "予算・選定を広範に委譲", "KPI設定・ガバナンス", "将来オプション（要実証）"],
    ]
    add_table(slide, data, Inches(0.45), Inches(1.25), Inches(12.40), Inches(2.55))
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(4.05), Inches(12.33), Inches(1.35))
    note.fill.solid()
    note.fill.fore_color.rgb = IBM_LIGHT
    note.line.color.rgb = ACCENT
    note.text_frame.text = (
        "IBM推奨（Option Bとの対応）：トラックB実証は Level 1 から開始 → 6ヶ月の実証後に Level 2 への拡大可否を判断\n"
        "経営層ディスカッション：Level 1 の具体的範囲（ツールカタログ・予算上限・期間）を確定することが最優先\n"
        "※ Level定義は論点整理用の案です。最終的な委譲範囲はクライアント様ご判断により確定します。"
    )
    for p in note.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT
    add_footer(slide, "権限成熟度モデル", slide_no)


def build_concept_slide(slide, slide_no):
    set_bg(slide)
    add_label(slide, "第2部｜説明の骨格")
    add_title(slide, "発電所自立経営とは何か｜Why / What")
    why = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(1.25), Inches(6.00), Inches(2.35))
    why.fill.solid()
    why.fill.fore_color.rgb = IBM_LIGHT
    why.line.color.rgb = IBM_BLUE
    add_bullets(slide, Inches(0.70), Inches(1.35), Inches(5.60), Inches(2.15), [
        ("Why｜なぜ今、定義が必要か", True),
        "・4月より発電所長が経営責任者に移行済みだが、「何を自律的に運営するか」の合意が未整理",
        "・権限委譲（現場裁量）とKPI管理（本社集中）の非対称が、施策実行のボトルネック",
        "・経営層で「どこまで委譲するか」を議論する前提として、概念の共通理解が必要",
    ])
    what = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.70), Inches(1.25), Inches(6.13), Inches(2.35))
    what.fill.solid()
    what.fill.fore_color.rgb = GREEN_BG
    what.line.color.rgb = GREEN_LINE
    add_bullets(slide, Inches(6.90), Inches(1.35), Inches(5.73), Inches(2.15), [
        ("What｜発電所自立経営とは（クライアント文脈）", True),
        "・発電所長がKPI達成に向け、権限の範囲内で自律的に運営・施策実行するモデル",
        "・アセット管理構築プログラム：収支計画PDCA・ユニット価値最大化　／　DX：デジタル施策の実行手段",
        "・≠ トラックA（AI検証）（技術・導入検証）　／　＝ 権限×KPI×PDCAの運用モデル実証（トラックB）",
    ])
    dist = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(3.85), Inches(12.33), Inches(1.55))
    dist.fill.solid()
    dist.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    dist.line.color.rgb = ORANGE
    add_bullets(slide, Inches(0.70), Inches(3.95), Inches(11.90), Inches(1.35), [
        ("So What｜本資料の提案との関係", True),
        "・二層体制（KPI管理ライン／実行権限ライン）で「自律」と「統制」を両立",
        "・権限成熟度 Level 1 から段階的に委譲（第1部参照）／ トラックB モデルプラントで実証",
        "・次スライド以降：3つの取組の関係 → ギャップ → 体制案 → Options の順で具体化",
    ])
    add_footer(slide, "発電所自立経営の定義", slide_no)


def replace_text_in_slide(slide, mapping):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    for old, new in mapping.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            for old, new in mapping.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)


def set_slide_number(slide, num):
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t.isdigit() and shape.top.inches > 7.0 and shape.width.inches < 0.5:
                shape.text_frame.text = str(num)


def main():
    if not SRC.exists():
        sys.exit(
            f"Input deck not found: {SRC}\n"
            "Place v5.1 at outputs/plant-autonomy_operating-model_draft_v5.1.pptx first."
        )
    src = Presentation(str(SRC))
    dst = Presentation()
    dst.slide_width = src.slide_width
    dst.slide_height = src.slide_height

    # v5.1 indices 0-based: clone order with inserts
    # 0-3: slides 1-4, NEW maturity, 4: slide 5 schedule, NEW concept, 5-18: slides 6-19
    plan = [
        ("clone", 0),
        ("clone", 1),
        ("clone", 2),
        ("clone", 3),
        ("new_maturity",),
        ("clone", 4),
        ("new_concept",),
        ("clone", 5),
        ("clone", 6),
        ("clone", 7),
        ("clone", 8),
        ("clone", 9),
        ("clone", 10),
        ("clone", 11),
        ("clone", 12),
        ("clone", 13),
        ("clone", 14),
        ("clone", 15),
        ("clone", 16),
        ("clone", 17),
        ("clone", 18),
    ]

    slide_no = 0
    for step in plan:
        slide_no += 1
        if step[0] == "clone":
            slide = clone_slide(src, step[1], dst)
        elif step[0] == "new_maturity":
            slide = dst.slides.add_slide(dst.slide_layouts[0])
            build_maturity_slide(slide, slide_no)
        elif step[0] == "new_concept":
            slide = dst.slides.add_slide(dst.slide_layouts[0])
            build_concept_slide(slide, slide_no)
        set_slide_number(slide, slide_no)

    repl = {
        "v5.1": "v5.2",
        "たたき台 v5.1": "たたき台 v5.2",
        "本v5.1を共有": "本v5.2を共有",
        "全3枚": "全4枚",
        "結論サマリー・経営判断3点・スケジュール概要": "結論サマリー・経営判断3点・権限成熟度・スケジュール概要",
        "全5枚": "全6枚",
        "全体プロセス・ギャップ整理・体制案・モデルプラント概要・Options": "自立経営の定義・全体プロセス・ギャップ整理・体制案・モデルプラント概要・Options",
        "全8枚（索引含め9枚）": "全8枚（索引含め10枚）",
    }
    for slide in dst.slides:
        replace_text_in_slide(slide, repl)

    # Page refs for v5.2 (explicit — avoid double-replace)
    page_refs = {
        12: (  # slide 13 index
            "詳細データ・質疑対応用\n"
            "4月移行後の実態ギャップ（詳細：プロセス・PJ連携）　p.14\n"
            "クライアント推進体制案への示唆（詳細）　p.15\n"
            "権限委譲マトリクス詳細・意思決定責任　p.16\n"
            "アセット管理構築プログラムとの接続　p.17\n"
            "モデルプラント検証設計・スケジュール・役割分担（詳細）　p.18-20\n"
            "経営判断事項・次のステップ　p.21"
        ),
        8: None,  # slide 9 gap summary footnote handled below
    }
    for idx, text in page_refs.items():
        if text is None:
            continue
        updated = False
        for shape in dst.slides[idx].shapes:
            if shape.has_text_frame and ("詳細データ" in shape.text_frame.text or "p.1" in shape.text_frame.text):
                if not updated:
                    shape.text_frame.text = text
                    updated = True
                else:
                    shape.text_frame.text = ""

    for shape in dst.slides[8].shapes:
        if shape.has_text_frame and "p.1" in shape.text_frame.text:
            shape.text_frame.text = "※ プロセス・PJ連携の観点の詳細は第3部参考資料（p.14）を参照"

    # Slide 4 IBM note
    for shape in dst.slides[3].shapes:
        if shape.has_text_frame and "IBM推奨" in shape.text_frame.text:
            shape.text_frame.text = (
                "IBM推奨：段階的委譲（Option B）＋PMO機能拡張＋トラックBはトラックA開始後に段階選定。"
                "委譲Level1の範囲確定を最優先の判断事項とする"
                "（Level定義は次スライド参照／詳細Optionsは第2部参照）"
            )

    for out in OUT_PATHS:
        out.parent.mkdir(parents=True, exist_ok=True)
        dst.save(str(out))
        print(f"Saved: {out}")
    print(f"Total slides: {len(dst.slides)}")


if __name__ == "__main__":
    main()
